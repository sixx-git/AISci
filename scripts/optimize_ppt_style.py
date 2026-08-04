"""
PPT Style Optimizer for AISci 互联网+ Competition (v6 → v7)
==========================================================
Based on analysis of:
  - 国内首创足球技战术实时分析平台PPT.pdf  (football reference: gold art text, nav bar, accent bars)
  - 航运天地一体通信服务保障系统PPT.pptx   (shipping reference: pill labels, ribbon nav, stat cards)
  - Web research on 互联网+ award-winning PPT design patterns

Optimizations applied:
  1. Gradient art-text (WordArt style) on all major titles via XML <a:gradFill> + glow + shadow
  2. Vertical cyan gradient accent bar before section titles (S02-S18)
  3. Enhanced cover page: main title + subtitle as gradient art text
  4. Enhanced ending page: slogan as gradient art text
"""

import os
import shutil
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from lxml import etree

# ── Namespaces ──────────────────────────────────────────────
A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

# ── Design Constants ─────────────────────────────────────────
GRAD_CYAN_BLUE = ['00E5FF', '3B82F6']   # Cyan → Blue (AI/Tech feel)
GRAD_GOLD      = ['FFE066', 'FFB300']    # Warm Gold (award feel)
GLOW_CYAN      = '00E5FF'
GLOW_GOLD      = 'FFD54A'

# Section name mapping for S02-S18
SECTIONS = {
    2: '项目价值', 3: '政策背景', 4: '痛点分析', 5: '解决方案',
    6: '产品体系', 7: '产品Demo', 8: '技术架构', 9: '核心创新',
    10: '技术壁垒', 11: '竞品分析', 12: '商业模式', 13: '市场推广',
    14: '应用成果', 15: '团队介绍', 16: '教育与社会价值',
    17: '财务与融资', 18: '未来规划',
}


# ── XML Builders ────────────────────────────────────────────

def _ns(tag):
    """Return fully-qualified tag name for DrawingML namespace."""
    return '{%s}%s' % (A, tag)


def _make_gradient_fill(parent, colors, angle_deg=90):
    """Create <a:gradFill> with linear gradient."""
    gf = etree.SubElement(parent, _ns('gradFill'))
    gsl = etree.SubElement(gf, _ns('gsLst'))
    n = len(colors)
    for i, c in enumerate(colors):
        pos = int(i * 100000 / max(n - 1, 1))
        gs = etree.SubElement(gsl, _ns('gs'))
        gs.set('pos', str(pos))
        sc = etree.SubElement(gs, _ns('srgbClr'))
        sc.set('val', c)
    ln = etree.SubElement(gf, _ns('lin'))
    ln.set('ang', str(angle_deg * 60000))   # 60k EMU per degree
    ln.set('scaled', '0')
    return gf


def _make_glow(parent, color, rad=40000, alpha_val=50000):
    """Add <a:glow> to an effectLst element."""
    gl = etree.SubElement(parent, _ns('glow'))
    gl.set('rad', str(rad))
    sc = etree.SubElement(gl, _ns('srgbClr'))
    sc.set('val', color)
    al = etree.SubElement(sc, _ns('alpha'))
    al.set('val', str(alpha_val))


def _make_shadow(parent, blur=50800, dist=25400,
                 dir_=2700000, alpha_val=45000):
    """Add <a:outerShdw> (soft drop shadow)."""
    sh = etree.SubElement(parent, _ns('outerShdw'))
    sh.set('blurRad', str(blur))
    sh.set('dist', str(dist))
    sh.set('dir', str(dir_))
    sh.set('algn', 'tl')
    sc = etree.SubElement(sh, _ns('srgbClr'))
    sc.set('val', '000000')
    al = etree.SubElement(sc, _ns('alpha'))
    al.set('val', str(alpha_val))


# ── Core: Apply Art-Text Effect to a Run ───────────────────

def apply_art_text(run, colors, angle=90, glow_color=None,
                   add_shadow=True, force_bold=True):
    """
    Replace run's solid fill with gradient + optional glow + shadow.
    This creates WordArt-style text that renders beautifully in PowerPoint.
    """
    rpr = run._r.get_or_add_rPr()

    # Remove existing fill elements (solidFill / gradFill)
    for old in list(rpr):
        tag = old.tag.split('}')[-1] if '}' in old.tag else old.tag
        if tag in ('solidFill', 'gradFill'):
            rpr.remove(old)

    if force_bold:
        rpr.set('b', '1')

    # Gradient fill
    _make_gradient_fill(rpr, colors, angle)

    # Effect list (glow + shadow)
    eff = etree.SubElement(rpr, _ns('effectLst'))
    if glow_color:
        _make_glow(eff, glow_color)
    if add_shadow:
        _make_shadow(eff)


# ── Per-Slide Processors ───────────────────────────────────

def process_cover(slide):
    """S01: Apply gradient art-text to main title + subtitle."""
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        top_in = sh.top / 914400
        wid_in = sh.width / 914400

        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                txt = run.text.strip()
                if not txt:
                    continue
                sz = run.font.size

                # Main title "联邦智研 AISci" — ~56pt at y≈2.5in
                if abs(top_in - 2.5) < 0.4 and wid_in > 5 and sz and sz >= Pt(40):
                    apply_art_text(run, GRAD_CYAN_BLUE, angle=90,
                                   glow_color=GLOW_CYAN, add_shadow=True)
                    print(f"    ★ Cover main title art-text: {txt[:30]}")

                # Subtitle "下一代 AI 科研生产力平台" — ~26pt at y≈4.0in
                elif abs(top_in - 4.0) < 0.4 and wid_in > 5 and sz and sz >= Pt(20):
                    apply_art_text(run, GRAD_GOLD, angle=0,
                                   glow_color=GLOW_GOLD, add_shadow=False)
                    print(f"    ★ Cover subtitle art-text: {txt[:30]}")


def process_ending(slide):
    """S19: Apply gradient art-text to main slogan."""
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        top_in = sh.top / 914400
        wid_in = sh.width / 914400

        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                txt = run.text.strip()
                if not txt or len(txt) < 8:
                    continue
                sz = run.font.size

                # Main slogan — ~38pt at y≈2.5in, full width
                if abs(top_in - 2.5) < 0.4 and wid_in > 10 and sz and sz >= Pt(28):
                    apply_art_text(run, GRAD_CYAN_BLUE, angle=90,
                                   glow_color=GLOW_CYAN, add_shadow=True)
                    print(f"    ★ Ending slogan art-text: {txt[:40]}")


def process_content_slide(slide, sidx):
    """S02-S18: Art-text on section title + vertical accent bar."""
    sec_name = SECTIONS.get(sidx)

    # 1) Find & style the section title run
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        top_in = sh.top / 914400
        left_in = sh.left / 914400
        wid_in = sh.width / 914400

        # Section title box: (0.3, 0.2), w<4in, contains the chapter name
        is_title_box = (
            abs(top_in - 0.2) < 0.15
            and abs(left_in - 0.3) < 0.2
            and wid_in < 4.0
        )

        if is_title_box:
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    txt = run.text.strip()
                    if not txt:
                        continue
                    if sec_name and (txt == sec_name or sec_name in txt):
                        apply_art_text(
                            run, GRAD_CYAN_BLUE, angle=90,
                            glow_color=GLOW_CYAN, add_shadow=False
                        )
                        print(f"    ★ Section title art-text: {txt}")

    # 2) Add cyan accent bar left of title
    if sec_name:
        _add_accent_bar(slide)


def _add_accent_bar(slide):
    """
    Add a thin vertical gradient bar at (0.05, 0.18).
    Positioned just left of the section title (which starts at x≈0.3).
    """
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(int(0.05 * 914400)),   # left  = 0.05 in
        Emu(int(0.18 * 914400)),   # top   = 0.18 in
        Emu(int(0.07 * 914400)),   # width = 0.07 in (thin bar)
        Emu(int(0.64 * 914400)),   # height = 0.64 in
    )
    bar.line.fill.background()       # no border

    # Cyan → Blue gradient fill
    fill = bar.fill
    fill.gradient()
    fill.gradient_angle = 90
    fill.gradient_stops[0].color.rgb = RGBColor(0x00, 0xE5, 0xFF)
    fill.gradient_stops[1].color.rgb = RGBColor(0x3B, 0x82, 0xF6)

    # Send to back (behind text but above background image)
    sp_tree = slide.shapes._spTree
    elem = bar._element
    sp_tree.remove(elem)
    # Insert after position 2 (typically bg picture + scrim shape)
    sp_tree.insert(2, elem)


# ── Main ───────────────────────────────────────────────────

def main():
    base = r'D:\Workplace\AISci\output'
    src = os.path.join(base, 'AISci_互联网+答辩PPT_v6.pptx')
    dst = os.path.join(base, 'AISci_互联网+答辩PPT_v7.pptx')

    if not os.path.exists(src):
        raise FileNotFoundError(f"Source not found: {src}")

    shutil.copy2(src, dst)
    print(f"Copied v6 → v7 ({os.path.getsize(dst)/1e6:.1f} MB)")

    prs = Presentation(dst)
    total = len(prs.slides)
    print(f"Processing {total} slides...\n")

    for i, slide in enumerate(prs.slides, 1):
        if i == 1:
            print(f"S{i:02d} [Cover]:")
            process_cover(slide)
        elif i == 19:
            print(f"S{i:02d} [Ending]:")
            process_ending(slide)
        elif i in SECTIONS:
            print(f"S{i:02d} [{SECTIONS[i]}]:")
            process_content_slide(slide, i)
        else:
            print(f"S{i:02d}: (skip)")

    prs.save(dst)
    size_mb = os.path.getsize(dst) / 1e6
    print(f"\n{'='*50}")
    print(f"✅ Done! Saved → {os.path.basename(dst)} ({size_mb:.1f} MB)")
    print(f"   Art-text titles: cover + ending + {len(SECTIONS)} content slides")
    print(f"   Accent bars:     {len(SECTIONS)} content slides")


if __name__ == '__main__':
    main()
