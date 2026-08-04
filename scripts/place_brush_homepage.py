"""
将毛笔书法"联邦智研"放到首页，替换原主标题艺术字图。
同时渲染"AISci"小字标放在毛笔字右侧。
输出 v10.pptx。
"""
import os, sys, shutil
from PIL import Image, ImageDraw, ImageFont, ImageFilter
import numpy as np

# ── 配置 ──
V9_PATH = r'D:\Workplace\AISci\output\AISci_互联网+答辩PPT_v9.pptx'
V10_PATH = r'D:\Workplace\AISci\output\AISci_互联网+答辩PPT_v10.pptx'
BRUSH_PNG = r'D:\Workplace\AISci\output\_brush\lianbang_zhiyan_brush.png'
ART_DIR   = r'D:\Workplace\AISci\output\_art_text'

# ── 1. 渲染 "AISci" 小字标（金铜渐变 + 描边）──
def render_aisci_tag(out_path, text='AISci', font_size=36,
                     color_top='#F5DEB3', color_bot='#B8860B',
                     stroke_color='#8B6914', stroke_width=2):
    font_path = r'C:\Windows\Fonts\HarmonyOS_Sans_SC_Bold.ttf'
    if not os.path.exists(font_path):
        font_path = r'C:\Windows\Fonts\msyhbd.ttc'
    font = ImageFont.truetype(font_path, font_size)

    # 测量尺寸
    dummy = Image.new('RGBA', (1, 1))
    draw = ImageDraw.Draw(dummy)
    bbox = draw.textbbox((0, 0), text, font=font)
    tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
    pad = font_size // 2
    img_w, img_h = tw + pad * 2, th + pad * 2

    img = Image.new('RGBA', (img_w, img_h), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)

    # 描边
    stroke_rgb = tuple(int(stroke_color[i:i+2], 16) for i in (1, 3, 5))
    for dx in range(-stroke_width, stroke_width + 1):
        for dy in range(-stroke_width, stroke_width + 1):
            if dx * dx + dy * dy <= stroke_width * stroke_width:
                draw.text((pad + dx, pad + dy), text, fill=(*stroke_rgb, 180), font=font)

    # 渐变填充
    top_rgb = np.array([int(color_top[i:i+2], 16) for i in (1, 3, 5)], dtype=np.float64)
    bot_rgb = np.array([int(color_bot[i:i+2], 16) for i in (1, 3, 5)], dtype=np.float64)

    # 逐像素渲染渐变文字（用文字mask）
    txt_img = Image.new('L', (img_w, img_h), 0)
    txt_draw = ImageDraw.Draw(txt_img)
    txt_draw.text((pad, pad), text, fill=255, font=font)
    txt_arr = np.array(txt_img, dtype=np.float64) / 255.0

    rgba = np.zeros((img_h, img_w, 4), dtype=np.uint8)
    for y in range(img_h):
        t = y / max(img_h - 1, 1)
        c = (top_rgb * (1 - t) + bot_rgb * t).astype(np.uint8)
        mask = txt_arr[y, :]
        rgba[y, :, 0] = c[0] * mask
        rgba[y, :, 1] = c[1] * mask
        rgba[y, :, 2] = c[2] * mask
        rgba[y, :, 3] = (mask * 230).astype(np.uint8)  # 稍微半透融合

    result = Image.fromarray(rgba, mode='RGBA')
    result.save(out_path, optimize=True)
    print(f'  AISci tag: {out_path} ({tw}x{th} -> {img_w}x{img_h})')
    return out_path


# ── 2. 主流程 ──
def main():
    from pptx import Presentation
    from pptx.util import Emu, Inches
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from lxml import etree

    NS_B = '{http://schemas.openxmlformats.org/drawingml/2006/main}blip'
    NS_R = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'

    # 渲染 AISci 标签
    aisci_path = os.path.join(ART_DIR, 'aisci_tag.png')
    render_aisci_tag(aisci_path)

    # 复制 v9 -> v10
    shutil.copy2(V9_PATH, V10_PATH)
    print(f'\nCopied v9 -> v10 ({os.path.getsize(V10_PATH)/1e6:.1f} MB)')

    prs = Presentation(V10_PATH)
    slide = prs.slides[0]  # 首页

    # 读取毛笔图尺寸
    brush_im = Image.open(BRUSH_PNG)
    bw, bh = brush_im.size
    print(f'Brush image: {bw}x{bh} (ratio={bw/bh:.2f})')

    # 定位：放在原主标题区域附近
    # 原 Picture 19: (0.50, 2.20) size=8.10x1.80
    # 毛笔字横排4字，希望更宽一点、稍微高一点以显气势
    target_w_inch = 8.5   # 英寸宽
    target_h_inch = target_w_inch / (bw / bh)  # 保持比例
    print(f'Target size: {target_w_inch:.1f} x {target_h_inch:.1f} inch')

    # 如果太高则限制高度，允许轻微压缩
    MAX_H = 3.2
    if target_h_inch > MAX_H:
        target_h_inch = MAX_H
        target_w_inch = target_h_inch * (bw / bh)
        print(f'Capped height: {target_w_inch:.1f} x {target_h_inch:.1f} inch')

    # 放置位置：左对齐原标题区，略偏上
    brush_left = Inches(0.55)
    brush_top = Inches(1.45)

    # ── 插入毛笔字图片 ──
    brush_pic = slide.shapes.add_picture(
        BRUSH_PNG, brush_left, brush_top,
        width=Inches(target_w_inch), height=Inches(target_h_inch)
    )
    print(f'Inserted brush at ({0.55:.2f}, {1.45:.2f}) size={target_w_inch:.1f}x{target_h_inch:.1f}')

    # ── 插入 AISci 小标签（毛笔字右下方）──
    aisci_left = brush_left + Inches(target_w_inch) - Inches(2.8)  # 靠右对齐
    aisci_top = brush_top + Inches(target_h_inch) - Inches(0.6)   # 底部偏上
    aisci_pic = slide.shapes.add_picture(
        aisci_path, aisci_left, aisci_top
    )
    print(f'Inserted AISci tag at ({aisci_left/914400:.2f}, {aisci_top/914400:.2f})')

    # ── 删除原主标题艺术字图 (Picture 19) ──
    removed_name = None
    for sh in list(slide.shapes):
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            left = (sh.left or 0) / 914400
            top = (sh.top or 0) / 914400
            w = (sh.width or 0) / 914400
            # Picture 19: (0.50, 2.20) 8.10x1.80
            if abs(left - 0.50) < 0.1 and abs(top - 2.20) < 0.3 and w > 7:
                sp = sh._element.getparent()
                sp.remove(sh._element)
                removed_name = sh.name
                print(f'Removed old title: {removed_name}')
                break

    # ── 保存 ──
    prs.save(V10_PATH)
    size_mb = os.path.getsize(V10_PATH) / 1e6
    print(f'\n✅ Saved: {V10_PATH} ({size_mb:.1f} MB)')


if __name__ == '__main__':
    main()
