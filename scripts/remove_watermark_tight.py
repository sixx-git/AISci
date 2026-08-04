"""
精确紧贴掩码去水印 v2 —— 基于 v4 背景图，仅覆盖即梦水印实际区域。
修复: 正确的 blip 查找+替换; 防止大范围误检(面积上限); 默认兜底掩码。
"""
import os, sys, time, shutil, io
import numpy as np
from PIL import Image
from scipy import sparse
from scipy.sparse.linalg import spsolve
from scipy.ndimage import binary_dilation, binary_closing, label
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

NS_B = "{http://schemas.openxmlformats.org/drawingml/2006/main}blip"
NS_R = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"


def build_laplacian_and_solve(img_arr, mask):
    h, w, c = img_arr.shape
    result = img_arr.copy()
    ys, xs = np.where(mask)
    n = len(ys)
    if n == 0:
        return result, 0, 0
    idx = np.full((h, w), -1, dtype=np.int32)
    idx[ys, xs] = np.arange(n)
    row_list, col_list, data_list = [], [], []
    b = np.zeros((n, c))
    for k in range(n):
        y, x = ys[k], xs[k]
        row_list.append(k); col_list.append(k); data_list.append(4.0)
        bnd_sum = np.zeros(c)
        for dy, dx in [(-1,0),(1,0),(0,-1),(0,1)]:
            ny, nx_ = y+dy, x+dx
            if 0 <= ny < h and 0 <= nx_ < w:
                if mask[ny, nx_]:
                    ni = idx[ny, nx_]
                    row_list.append(k); col_list.append(ni); data_list.append(-1.0)
                else:
                    bnd_sum += img_arr[ny, nx_, :]
        b[k] = bnd_sum
    A = sparse.csr_matrix((data_list, (row_list, col_list)), shape=(n, n))
    t0 = time.time()
    for ch in range(c):
        result[ys, xs, ch] = np.clip(spsolve(A, b[:, ch]), 0, 255)
    return result, n, time.time()-t0


def detect_watermark_bbox(img_array,
                          search_x_min=0.72,
                          search_y_min=0.78,
                          lum_thresh=130,
                          min_area_px=200,
                          max_area_px=12000):
    h, w = img_array.shape[:2]
    gray = np.mean(img_array.astype(np.float64), axis=2)
    sx1 = int(w * search_x_min)
    sy1 = int(h * search_y_min)
    roi = gray[sy1:h, sx1:w]
    bright_mask = roi > lum_thresh
    kernel = np.ones((5, 5), dtype=bool)
    bright_mask = binary_closing(bright_mask, structure=kernel, iterations=3)
    labeled, num_features = label(bright_mask)
    if num_features == 0:
        return None
    sizes = [(labeled == i).sum() for i in range(1, num_features + 1)]
    max_label = np.argmax(sizes) + 1
    wm_mask = (labeled == max_label)
    area = wm_mask.sum()
    if area < min_area_px or area > max_area_px:
        print(f"    面积异常({area}px), 跳过自动检测")
        return None
    coords = np.where(wm_mask)
    y1 = sy1 + int(coords[0].min())
    y2 = sy1 + int(coords[0].max()) + 1
    x1 = sx1 + int(coords[1].min())
    x2 = sx1 + int(coords[1].max()) + 1
    pct = (x2-x1)*(y2-y1)/(h*w)*100
    print(f"    水印: ({x1},{y1})-({x2},{y2}) {x2-x1}x{y2-y1} {area}px {pct:.2f}%")
    return (y1, y2, x1, x2)


def make_tight_mask(h, w, bbox, margin=12):
    if bbox is None:
        # 兜底：右下角小默认掩码（仅覆盖典型水印位置）
        mask = np.zeros((h, w), dtype=bool)
        y1 = int(h * 0.87)
        x1 = int(w * 0.80)
        mask[y1:, x1:] = True
        kernel = np.ones((8, 8), dtype=bool)
        mask = binary_dilation(mask, structure=kernel, iterations=1)
        print(f"    使用默认掩码: {mask.sum()}px ({mask.sum()/(h*w)*100:.2f}%)")
        return mask
    y1, y2, x1, x2 = bbox
    y1 = max(0, y1 - margin)
    y2 = min(h, y2 + margin)
    x1 = max(0, x1 - margin)
    x2 = min(w, x2 + margin)
    mask = np.zeros((h, w), dtype=bool)
    mask[y1:y2, x1:x2] = True
    kernel = np.ones((6, 6), dtype=bool)
    mask = binary_dilation(mask, structure=kernel, iterations=1)
    print(f"    紧贴掩码: {mask.sum()}px ({mask.sum()/(h*w)*100:.2f}%)")
    return mask


def process_one(src_path, dst_path=None):
    if dst_path is None:
        dst_path = src_path
    im = Image.open(src_path).convert('RGB')
    arr = np.array(im, dtype=np.float64)
    h, w = arr.shape[:2]
    name = os.path.basename(src_path)[:50]
    print(f"\n  [{name}] {w}x{h}")
    bbox = detect_watermark_bbox(arr)
    mask = make_tight_mask(h, w, bbox, margin=12)
    result, n_px, elapsed = build_laplacian_and_solve(arr, mask)
    out = Image.fromarray(np.clip(result, 0, 255).astype(np.uint8))
    out.save(dst_path, optimize=True)
    check_gray = np.array(out, dtype=np.float64).mean(axis=2)
    remain = ((check_gray > 170) & mask).sum()
    status = "OK" if remain < 50 else f"WARN({remain})"
    print(f"    完成: {n_px}px/{elapsed:.2f}s 残留={remain} {status}")
    return True


def replace_bg_blob(slide, new_img_data):
    """用新图片数据替换幻灯片全出血背景图的 blob（直接替换已有 part 的字节）"""
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE and sh.left is not None and abs(sh.left) < 10000:
            wid = (sh.width or 0) / 914400
            if wid > 12:
                blip = sh._element.find(".//" + NS_B)
                if blip is not None:
                    rid = blip.get(NS_R)
                    if rid:
                        embed_rel = sh.part.rels.get(rid)
                        if embed_rel and embed_rel.target_part:
                            embed_rel.target_part._blob = new_img_data
                            return True
    return False


def main():
    base = r"D:\Workplace\AISci\output"
    pptx_v4 = os.path.join(base, "AISci_互联网+答辩PPT_v4.pptx")

    # Step 1: 提取背景图
    extract_dir = os.path.join(base, "_tight_wm_extract")
    os.makedirs(extract_dir, exist_ok=True)

    prs = Presentation(pptx_v4)
    extracted = []
    for sidx, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE and sh.left is not None and abs(sh.left) < 10000:
                wid = (sh.width or 0) / 914400
                if wid > 12:
                    outp = os.path.join(extract_dir, f"bg_{sidx:02d}.{sh.image.ext}")
                    with open(outp, "wb") as f:
                        f.write(sh.image.blob)
                    extracted.append((sidx, outp))
                    break

    print(f"提取 {len(extracted)} 张背景图")

    # Step 2: 精确去水印
    processed_dir = os.path.join(base, "_tight_wm_processed")
    os.makedirs(processed_dir, exist_ok=True)

    results = []
    for sidx, src_path in extracted:
        dst_path = os.path.join(processed_dir, f"bg_{sidx:02d}.png")
        process_one(src_path, dst_path)
        results.append((sidx, dst_path))

    # Step 3: 复制 v4 -> v6, 替换 blob
    out_pptx = os.path.join(base, "AISci_互联网+答辩PPT_v6.pptx")
    shutil.copy2(pptx_v4, out_pptx)
    print(f"\n复制 v4 -> v6")

    prs2 = Presentation(out_pptx)
    replaced = 0
    for sidx, dst_path in results:
        slide = prs2.slides[sidx - 1]
        with open(dst_path, "rb") as f:
            img_data = f.read()
        if replace_bg_blob(slide, img_data):
            replaced += 1
            print(f"  OK p{sidx:02d}")
        else:
            print(f"  FAIL p{sidx:02d}")

    prs2.save(out_pptx)
    size_mb = os.path.getsize(out_pptx) / 1e6
    print(f"\n{'='*60}")
    print(f"完成! 替换 {replaced}/{len(results)} 张 | {size_mb:.1f} MB")


if __name__ == "__main__":
    main()
