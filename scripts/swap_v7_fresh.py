"""
v7: 从头构建 PPT —— 基于 v6 谐波修复后的干净源图。
关键改进：每次打开图片前先读文件字节确认是最新版本，避免缓存问题。
"""
import os, io, time
from PIL import Image
from pptx import Presentation
from pptx.util import Emu

# 用 v4（有 scrim + 无 vN）作为模板，只替换背景图
TEMPLATE = r"D:\Workplace\AISci\output\AISci_互联网+答辩PPT_v4.pptx"
OUTPUT   = r"D:\Workplace\AISci\output\AISci_互联网+答辩PPT_v5.pptx"

# 页面 → 图片文件夹映射 (1-indexed)
PAGE_MAP = {
    1: "第一页", 2: "第二页", 3: "第三页", 4: "第四页", 5: "第五页",
    6: "第六页", 7: "第七页", 8: "第八页", 9: "第九页", 10: "第十页",
    11: "第十一页", 12: "第十二页", 13: "第十三页", 14: "第十四页",
    15: "第十五页", 16: "第十六页", 17: "第十七页", 18: "第十八页",
    19: "第十九页",
}
IMG_BASE = r"D:\Workplace\AISci\output\图片"
MAX_W = 2000  # 缩放上限


def find_png(folder):
    """在文件夹中找唯一的 PNG"""
    d = os.path.join(IMG_BASE, folder)
    for f in os.listdir(d):
        if f.lower().endswith('.png'):
            return os.path.join(d, f)
    return None


def load_fresh_image(path):
    """强制从磁盘读取最新版本，缩放到 MAX_W"""
    # 先读原始字节确保拿到磁盘上的最新数据
    with open(path, 'rb') as f:
        raw = f.read()
    im = Image.open(io.BytesIO(raw)).convert('RGB')
    w, h = im.size
    if w > MAX_W:
        ratio = MAX_W / w
        im = im.resize((MAX_W, int(h * ratio)), Image.LANCZOS)
    buf = io.BytesIO()
    im.save(buf, format='PNG')
    buf.seek(0)
    print(f"    加载 {os.path.basename(path)[:40]} -> {im.size[0]}x{im.size[1]} ({len(raw)//1024}KB raw)")
    return buf


def main():
    t0 = time.time()
    prs = Presentation(TEMPLATE)
    SW, SH = prs.slide_width, prs.slide_height

    NS_B = '{http://schemas.openxmlformats.org/drawingml/2006/main}blip'
    NS_R = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'

    replaced = 0
    for slide_idx, folder in PAGE_MAP.items():
        slide = prs.slides[slide_idx - 1]
        png_path = find_png(folder)
        if not png_path:
            print(f"  ⚠️ Slide {slide_idx} ({folder}): 未找到图片")
            continue

        # 找到全出血背景图 shape
        bg_shape = None
        for sh in slide.shapes:
            if sh.shape_type == 13 and sh.left is not None:
                wid = sh.width / 914400 if sh.width else 0
                if wid > 12:  # 全出血
                    bg_shape = sh
                    break

        if bg_shape is None:
            print(f"  ⚠️ Slide {slide_idx}: 未找到背景图")
            continue

        # 加载新图片（强制从磁盘）
        img_buf = load_fresh_image(png_path)

        # 替换嵌入的媒体
        blip = bg_shape._element.find('.//' + NS_B)
        rid = blip.get(NS_R)
        rel = slide.part.rels[rid]

        # 删除旧关系和媒体，添加新的
        old_part = rel.target_part
        new_img_part = slide.part.package.add_image(img_buf.getvalue())
        slide.part.rels[rid]._target = new_img_part

        # 更新 blip 引用
        new_rid = slide.part.relate_to(new_img_part, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/image')
        blip.set(NS_R, new_rid)

        replaced += 1
        print(f"  ✅ Slide {slide_idx} ({folder}): 背景已替换")

    prs.save(OUTPUT)
    dt = time.time() - t0
    sz = os.path.getsize(OUTPUT) / 1e6
    print(f"\n{'='*50}")
    print(f"完成! 替换 {replaced}/{len(PAGE_MAP)} 张背景")
    print(f"输出: {OUTPUT}")
    print(f"大小: {sz:.1f} MB | 耗时: {dt:.1f}s")


if __name__ == '__main__':
    main()
