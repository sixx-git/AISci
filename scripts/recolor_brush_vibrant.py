"""
毛笔书法重新染色：
- 联邦 / 研 三个字：更鲜艳的金铜渐变
- 智 字：鲜艳青蓝渐变，羽化过渡避免硬边
输出透明 PNG，并替换 v11 首页毛笔图为 v12。
"""
import numpy as np
from PIL import Image
import os, glob
from pptx import Presentation
from pptx.util import Emu

BRUSH_DIR = r'D:\Workplace\AISci\output\_brush'
orig = glob.glob(os.path.join(BRUSH_DIR, 'Traditional_Chinese_brush_call_*.png'))[0]
CLEAN = os.path.join(BRUSH_DIR, 'lianbang_zhiyan_brush_clean.png')
OUT = os.path.join(BRUSH_DIR, 'lianbang_zhiyan_brush_vibrant.png')
PPTX_IN = r'D:\Workplace\AISci\output\AISci_互联网+答辩PPT_v11.pptx'
PPTX_OUT = r'D:\Workplace\AISci\output\AISci_互联网+答辩PPT_v12.pptx'

# ── 1. 抠图 + 染色 ──
im = Image.open(orig).convert('RGB')
arr = np.array(im, dtype=np.float64)
h, w = arr.shape[:2]
lum = 0.299*arr[:,:,0] + 0.587*arr[:,:,1] + 0.114*arr[:,:,2]
a = np.clip((138-lum)/max(138-95, 1), 0, 1)
a[lum < 95] = 1.0
a[lum > 138] = 0.0
a[int(h*0.82):, int(w*0.78):] = 0.0   # 水印区透明

# 鲜艳金铜渐变 (base)
GOLD_TOP = np.array([255, 232, 115], float)  # #FFE873
GOLD_BOT = np.array([224, 138, 0], float)    # #E08A00
# 智字鲜艳青蓝渐变
CYAN_TOP = np.array([125, 249, 255], float)  # #7DF9FF
CYAN_BOT = np.array([30, 107, 255], float)   # #1E6BFF

y_norm = np.arange(h, dtype=float).reshape(-1, 1) / max(h-1, 1)

def grad(t, b):
    return (t[0]+(b[0]-t[0])*y_norm,
            t[1]+(b[1]-t[1])*y_norm,
            t[2]+(b[2]-t[2])*y_norm)

gr, gG, gb = grad(GOLD_TOP, GOLD_BOT)
cr, cG, cb = grad(CYAN_TOP, CYAN_BOT)

# 智字区域 (ink-mass 第3段 537..813) 带羽化
zlo, zhi = 537, 813
feather = 18
xx = np.arange(w)
mz = np.clip((xx-(zlo-feather))/feather, 0, 1) * np.clip(((zhi+feather)-xx)/feather, 0, 1)
mz = mz.reshape(1, -1)
mz = mz * (a > 0.02)   # 仅笔画处上色

# 保留墨色深浅 → 金属质感（提高下限让颜色更鲜艳）
ref = np.clip(lum/120.0, 0.45, 1.0)

r = (gr*(1-mz) + cr*mz) * ref
g = (gG*(1-mz) + cG*mz) * ref
b = (gb*(1-mz) + cb*mz) * ref

rgba = np.zeros((h, w, 4), np.uint8)
rgba[:,:,0] = np.clip(r, 0, 255).astype(np.uint8)
rgba[:,:,1] = np.clip(g, 0, 255).astype(np.uint8)
rgba[:,:,2] = np.clip(b, 0, 255).astype(np.uint8)
rgba[:,:,3] = (a*255).astype(np.uint8)
Image.fromarray(rgba, 'RGBA').save(OUT, optimize=True)
print('recolor ->', OUT, 'size', os.path.getsize(OUT)//1024, 'KB')
ink = (a > 0.05).sum()
print('  智字(青蓝)占笔画:', round(float((mz > 0.5).sum())/max(ink,1)*100, 1), '%')
print('  不透明率:', round(float(ink)/(h*w)*100, 1), '%')

# ── 2. 替换 v11 首页毛笔图 ──
new_blob = open(OUT, 'rb').read()
clean_blob = open(CLEAN, 'rb').read()
prs = Presentation(PPTX_IN)
sld = prs.slides[0]
replaced = 0
for sh in sld.shapes:
    if sh.shape_type == 13:   # PICTURE
        try:
            blob = sh.image.blob
        except Exception:
            continue
        if blob == clean_blob:
            blip = sh._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
            rid = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
            rel = sld.part.rels[rid]
            rel.target_part._blob = new_blob
            replaced += 1
            print('  replaced picture rid', rid, 'dims',
                  round(Emu(sh.width).inches, 2), 'x', round(Emu(sh.height).inches, 2))
prs.save(PPTX_OUT)
print('saved ->', PPTX_OUT, 'replaced', replaced)
