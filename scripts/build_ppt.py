#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Build AISci internet+ defense PPT (19 pages, dark tech theme).
Backgrounds + page visuals are PIL-generated (see gen_assets.py).
Text uses "等线" font; unified dark navy / electric blue / gold system.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

ASSET = "D:/Workplace/AISci/output/AISci_ppt/assets"
OUT = "D:/Workplace/AISci/output/AISci_互联网+答辩PPT.pptx"
TOTAL = 19

W = Inches(13.33); H = Inches(7.5)

# ---- palette ----
NAVY   = RGBColor(0x0A, 0x16, 0x28)
DARK   = RGBColor(0x0F, 0x1D, 0x35)
CARD   = RGBColor(0x14, 0x24, 0x42)
CARDL  = RGBColor(0x2C, 0x46, 0x74)
EBLUE  = RGBColor(0x38, 0x6F, 0xC4)
CYAN   = RGBColor(0x3F, 0xE0, 0xFF)
GOLD   = RGBColor(0xC9, 0xA8, 0x4C)
LGOLD  = RGBColor(0xE8, 0xD5, 0x8C)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xC7, 0xD2, 0xE0)
SUB    = RGBColor(0x8A, 0x9C, 0xB5)
RED    = RGBColor(0xE0, 0x6A, 0x6A)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]
FONT = "等线"

def set_ea(run, name=FONT):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)

def add_bg(slide, path):
    if os.path.exists(path):
        slide.background.fill.background()
        pic = slide.shapes.add_picture(path, Emu(0), Emu(0), W, H)
        sp = pic._element
        sp.getparent().remove(sp)
        slide.shapes._spTree.insert(2, sp)

def text_box(slide, l, t, w, h, text, size=18, color=LIGHT, bold=False,
             align=PP_ALIGN.LEFT, line_spacing=1.25, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(4)
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
    set_ea(r)
    if line_spacing: p.line_spacing = Pt(size*line_spacing)
    return tb

def multi_text(slide, l, t, w, h, lines, size=16, color=LIGHT, bold=False,
               align=PP_ALIGN.LEFT, line_spacing=1.3, space=8):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        c = color; b = bold
        if isinstance(line, tuple):
            line, c, b = line
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.color.rgb = c; r.font.bold = b
        set_ea(r)
        p.line_spacing = Pt(size*line_spacing)
    return tb

def gold_line(slide, l, t, w, h=0.045):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = GOLD; s.line.fill.background()
    return s

def badge(slide, text, l=0.3, t=0.22):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(2.7), Inches(0.62))
    s.fill.solid(); s.fill.fore_color.rgb = NAVY; s.line.fill.background()
    tf = s.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(22); r.font.color.rgb = WHITE; r.font.bold = True; set_ea(r)
    return s

def sub(slide, text, l=3.15, t=0.28):
    return text_box(slide, l, t, 9.8, 0.55, f"——{text}", size=17, color=LGOLD)

def card(slide, l, t, w, h, fill=CARD, line=CARDL):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = line; s.line.width = Pt(1.25)
    s.shadow.inherit = False
    return s

def right_visual(slide, name):
    p = os.path.join(ASSET, name)
    if os.path.exists(p):
        slide.shapes.add_picture(p, Inches(6.72), Inches(0.05), width=Inches(6.6))

def page_num(slide, n):
    text_box(slide, 12.05, 6.95, 1.1, 0.4, f"{n}/{TOTAL}", size=11, color=SUB, align=PP_ALIGN.RIGHT)

def content(slide, btext, stext, vname, n):
    add_bg(slide, os.path.join(ASSET, "content-bg.png"))
    badge(slide, btext)
    sub(slide, stext)
    right_visual(slide, vname)
    page_num(slide, n)
    # left text zone: x in [0.4, 6.4]

# ============ SLIDE 1: COVER ============
s = prs.slides.add_slide(BLANK)
add_bg(s, os.path.join(ASSET, "cover-bg.png"))
gold_line(s, 0.8, 2.35, 5.0)
text_box(s, 0.8, 0.5, 9.0, 0.5, "中国国际大学生创新大赛（2026）  ·  高教主赛道 / 人工智能与大数据", size=15, color=LGOLD, bold=True)
text_box(s, 0.8, 2.5, 7.5, 1.2, "联邦智研 AISci", size=56, color=WHITE, bold=True)
text_box(s, 0.8, 3.95, 7.5, 0.6, "下一代 AI 科研生产力平台", size=26, color=LGOLD, bold=True)
text_box(s, 0.8, 4.7, 6.5, 0.6, "让科研创新进入智能时代 —— 把『问答工具』升级为『可信的科研协作者』", size=15, color=LIGHT)
# info cards
for i, (k, v) in enumerate([("赛道", "高教主赛道 · 人工智能"), ("团队", "联邦智研团队"), ("负责人", "[请填姓名]")]):
    x = 0.8 + i*2.55
    card(s, x, 5.5, 2.35, 0.95, fill=RGBColor(0x13,0x22,0x3C), line=GOLD)
    text_box(s, x+0.15, 5.58, 2.05, 0.35, k, size=12, color=GOLD, bold=True)
    text_box(s, x+0.15, 5.9, 2.05, 0.45, v, size=13, color=WHITE)
right_visual(s, "v1_cover.png")

# ============ SLIDE 2: VALUE ============
s = prs.slides.add_slide(BLANK)
content(s, "项目价值", "用 AI 重构科研生产关系，让团队拥有『可信科研协作者』", "v2_value.png", 2)
card(s, 0.4, 1.5, 6.0, 2.5)
text_box(s, 0.65, 1.65, 5.5, 0.5, "传统科研流程（灰色调）", size=17, color=SUB, bold=True)
multi_text(s, 0.65, 2.2, 5.5, 1.6, [
    "文献搜索  →  人工分析  →  提出假设  →  实验验证",
    "耗时 · 易重复 · 难追溯",
    "（数据来源：商业计划书 v1 第一章）",
], size=15, color=LIGHT, line_spacing=1.4)
card(s, 0.4, 4.2, 6.0, 2.5, fill=RGBColor(0x10,0x22,0x44), line=EBLUE)
text_box(s, 0.65, 4.35, 5.5, 0.5, "AISci 智能科研闭环（蓝色调）", size=17, color=CYAN, bold=True)
multi_text(s, 0.65, 4.9, 5.5, 1.6, [
    "问题理解 → 证据检索 → 假设生成 → 验证优化 → 科研成果",
    "可审计 · 可复现 · 人机协同",
], size=15, color=LIGHT, line_spacing=1.4)
text_box(s, 0.4, 6.85, 6.2, 0.5, "通用 AI 是『问题→回答』；AISci 是『问题→证据→假设→验证→优化』的科研闭环。", size=13, color=LGOLD, bold=True)

# ============ SLIDE 3: POLICY ============
s = prs.slides.add_slide(BLANK)
content(s, "政策背景", "AI for Science 已上升为国家科技战略制高点", "v3_policy.png", 3)
card(s, 0.4, 1.45, 6.1, 5.3)
text_box(s, 0.65, 1.6, 5.6, 0.5, "三大战略支撑", size=18, color=WHITE, bold=True)
gold_line(s, 0.65, 2.15, 3.0)
multi_text(s, 0.65, 2.35, 5.6, 4.2, [
    ("① 国家《新一代人工智能发展规划》", EBLUE, True),
    "将『AI 赋能科学发现』列为重点方向，鼓励科研范式变革。",
    ("② 科技部『人工智能驱动的科学研究（AI4S）』专项", EBLUE, True),
    "支持从『人找知识』迈向『AI 协同发现』的范式跃迁。",
    ("③ 全球 AI for Science 市场高速增长", EBLUE, True),
    "权威机构预测年间复合增速领先，窗口期已经打开。",
], size=14.5, color=LIGHT, line_spacing=1.35, space=7)
text_box(s, 0.4, 6.85, 6.2, 0.5, "科研范式正从『人找知识』迈向『AI 协同发现』，这是一代人的窗口期。", size=13, color=LGOLD, bold=True)

# ============ SLIDE 4: PAIN ============
s = prs.slides.add_slide(BLANK)
content(s, "痛点分析", "信息爆炸与认知瓶颈，正在拖慢每一次科学发现", "v4_pain.png", 4)
pains = [
    ("信息过载", "全球年发文超 300 万篇，单点研究者难以全面追踪前沿。", "年发文 300万+"),
    ("创新盲区", "跨领域交叉处的机会，人工难以系统扫描，易重复已有工作。", "跨域机会难捕捉"),
    ("信任缺失", "通用 LLM 存在『幻觉引用 / 证据不可追溯 / 质量无控制』三大硬伤。", "幻觉·不可溯"),
]
for i, (title, desc, tag) in enumerate(pains):
    y = 1.5 + i*1.85
    card(s, 0.4, y, 6.0, 1.65)
    text_box(s, 0.65, y+0.12, 4.3, 0.5, title, size=19, color=WHITE, bold=True)
    text_box(s, 4.6, y+0.12, 1.7, 0.5, tag, size=13, color=GOLD, bold=True, align=PP_ALIGN.RIGHT)
    gold_line(s, 0.65, y+0.62, 2.2)
    text_box(s, 0.65, y+0.72, 5.6, 0.8, desc, size=14, color=LIGHT)

# ============ SLIDE 5: SOLUTION ============
s = prs.slides.add_slide(BLANK)
content(s, "解决方案", "把科研逻辑链编码为可审计的多智能体 Pipeline", "v5_radial.png", 5)
text_box(s, 0.4, 1.4, 6.1, 0.7, "科研智能体基础设施 + AI 科研生产力平台", size=19, color=WHITE, bold=True)
gold_line(s, 0.4, 2.1, 3.0)
anchors = [
    ("系统完整", "近产品原型（前端+后端+API+测试），非概念演示", EBLUE),
    ("国产适配", "基于 Qwen，对接信创与国产 AI 基础设施", CYAN),
    ("可信机制", "Evidence Chain + Boolean Gate + HITL 三位一体", GOLD),
    ("全流程闭环", "科学问题 → 可验证假设 → 实验迭代 → 报告生成", LGOLD),
]
for i, (t, d, c) in enumerate(anchors):
    y = 2.35 + i*1.12
    card(s, 0.4, y, 6.1, 1.0)
    text_box(s, 0.62, y+0.08, 2.0, 0.8, t, size=16, color=c, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, 2.6, y+0.08, 3.8, 0.85, d, size=13, color=LIGHT, anchor=MSO_ANCHOR.MIDDLE)
text_box(s, 0.4, 6.85, 6.2, 0.5, "我们刻意区别于『AI 科研助手』，做『可信的科研协作者』。", size=13, color=LGOLD, bold=True)

# ============ SLIDE 6: PRODUCT MATRIX ============
s = prs.slides.add_slide(BLANK)
content(s, "产品体系", "三条产品线，覆盖从个人科研到机构治理", "v6_matrix.png", 6)
prods = [
    ("AISci Research", "高校科研人员", ["AI 文献助手：语义检索·证据抽取·综述生成", "研究方向发现：知识缺口自动推荐", "科研计划生成：假设+方法+实验路线"]),
    ("AISci Enterprise", "企业研发", ["技术趋势分析 · 专利查新 / FTO", "创新方案生成 · 竞品技术追踪", "研发知识库与可审计报告"]),
    ("AISci Institution", "高校与科研机构", ["科研管理 · 项目评估（证据链评审）", "学科分析 · 团队能力画像", "科研数据治理与合规"]),
]
for i, (name, who, items) in enumerate(prods):
    y = 1.45 + i*1.82
    card(s, 0.4, y, 6.1, 1.7)
    text_box(s, 0.62, y+0.1, 4.0, 0.5, name, size=18, color=WHITE, bold=True)
    text_box(s, 4.5, y+0.12, 1.9, 0.5, who, size=12, color=GOLD, bold=True, align=PP_ALIGN.RIGHT)
    gold_line(s, 0.62, y+0.62, 2.0)
    multi_text(s, 0.62, y+0.72, 5.7, 0.9, [f"▸ {it}" for it in items], size=12.5, color=LIGHT, line_spacing=1.2, space=3)
text_box(s, 0.4, 6.85, 6.2, 0.5, "三条线共享同一套 Agent + Skill 底座，形成完整科研智能体生态。", size=13, color=LGOLD, bold=True)

# ============ SLIDE 7: DEMO ============
s = prs.slides.add_slide(BLANK)
content(s, "产品 Demo", "输入一个问题，输出一份可验证的科研方案", "v7_demo.png", 7)
text_box(s, 0.4, 1.4, 6.1, 0.5, "研究问题（示例）：如何降低新能源汽车电池热失控风险？", size=14, color=LGOLD, bold=True)
gold_line(s, 0.4, 1.95, 3.0)
steps = [
    ("Step 1  Research Agent", "分析已有论文，定位研究前沿"),
    ("Step 2  Knowledge Agent", "发现知识缺口，识别创新空间"),
    ("Step 3  Hypothesis Agent", "生成创新假设，给出研究思路"),
    ("Step 4  Evidence Agent", "验证证据链（Fact 白名单，无虚构引用）"),
    ("Step 5  Report Agent", "生成完整科研方案与可复现报告"),
]
for i, (t, d) in enumerate(steps):
    y = 2.1 + i*0.92
    card(s, 0.4, y, 6.1, 0.82)
    text_box(s, 0.6, y+0.06, 2.7, 0.7, t, size=14, color=EBLUE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, 3.3, y+0.06, 3.1, 0.7, d, size=12.5, color=LIGHT, anchor=MSO_ANCHOR.MIDDLE)
text_box(s, 0.4, 6.78, 6.2, 0.55, "右侧为系统运行示意，请替换为真实 Demo 截图。输入问题 → AI 科研闭环 → 输出成果。", size=12, color=SUB)

# ============ SLIDE 8: ARCHITECTURE ============
s = prs.slides.add_slide(BLANK)
content(s, "技术架构", "Agent—Skill—Infrastructure 三层分离，可私有化部署", "v8_arch.png", 8)
layers = [
    ("核心引擎层", "6 大智能体 + 七阶段 Pipeline", EBLUE),
    ("技能工具层", "8 大类、约 70 个可复用 Skill 模块", CYAN),
    ("基础设施层", "FastAPI（18 个 API 路由）+ React18 / Vite5 / TailwindCSS + SQLite + zvec 向量检索", GOLD),
]
for i, (t, d, c) in enumerate(layers):
    y = 1.5 + i*1.15
    card(s, 0.4, y, 6.1, 1.0)
    text_box(s, 0.62, y+0.08, 2.1, 0.8, t, size=16, color=c, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, 2.7, y+0.06, 3.7, 0.9, d, size=13, color=LIGHT, anchor=MSO_ANCHOR.MIDDLE)
text_box(s, 0.4, 5.05, 6.1, 1.9, "", size=12)
text_box(s, 0.4, 5.1, 6.1, 0.5, "七阶段 Pipeline（可 SaaS / 私有化多种交付）", size=14, color=WHITE, bold=True)
text_box(s, 0.4, 5.6, 6.1, 1.1, "问题理解 → 文献挖掘 → 知识缺口 → 假设生成 → 假设评审 → 迭代实验 → 报告生成\n每阶段自动记录审计数据，支持 jsonl 导出，完全可复现。", size=12.5, color=LIGHT, line_spacing=1.35)

# ============ SLIDE 9: INNOVATION ============
s = prs.slides.add_slide(BLANK)
content(s, "核心创新", "四项机制创新，最终都指向客户价值", "v9_innovation.png", 9)
inns = [
    ("证据链迭代推理 Evidence Chain", "Fact 白名单从机制杜绝幻觉引用", "商业价值：结果可信、可被客户审计", EBLUE),
    ("布尔质量门禁 Boolean Gate", "0/1 硬门禁替代连续评分", "商业价值：质量决策边界清晰、可停滞", CYAN),
    ("反馈中心 + 人在回路 HITL", "跨阶段传递约束 + 关键节点人工门控", "商业价值：自动化与可控性平衡", GOLD),
    ("反事实预演 Counterfactual", "L0 定性证伪过滤无效实验", "商业价值：保护科研资源、降本增效", LGOLD),
]
for i, (t, p, v, c) in enumerate(inns):
    y = 1.45 + i*1.32
    card(s, 0.4, y, 6.1, 1.2)
    text_box(s, 0.62, y+0.08, 5.7, 0.45, t, size=15, color=c, bold=True)
    text_box(s, 0.62, y+0.5, 5.7, 0.35, f"原理：{p}", size=12, color=LIGHT)
    text_box(s, 0.62, y+0.85, 5.7, 0.32, v, size=11.5, color=GOLD, bold=True)

# ============ SLIDE 10: MOAT ============
s = prs.slides.add_slide(BLANK)
content(s, "技术壁垒", "AISci 不是聊天机器人，而是科研智能体系统", "v10_barrier.png", 10)
card(s, 0.4, 1.5, 6.1, 2.3, fill=RGBColor(0x1A,0x1E,0x28), line=RGBColor(0x55,0x5E,0x70))
text_box(s, 0.62, 1.62, 5.6, 0.5, "ChatGPT（通用 LLM）", size=17, color=SUB, bold=True)
text_box(s, 0.62, 2.15, 5.6, 1.5, "回答问题\n（单次生成，无证据约束）", size=15, color=LIGHT, line_spacing=1.4)
card(s, 0.4, 4.0, 6.1, 2.6, fill=RGBColor(0x10,0x22,0x44), line=EBLUE)
text_box(s, 0.62, 4.12, 5.6, 0.5, "AISci（科研智能体）", size=17, color=CYAN, bold=True)
text_box(s, 0.62, 4.65, 5.6, 1.8, "理解问题 + 检索证据 + 生成假设\n+ 验证方案 + 持续优化", size=15, color=LIGHT, line_spacing=1.4)
text_box(s, 0.4, 6.78, 6.2, 0.55, "壁垒 = 流程完整性 + 可信机制 + 国产适配，组合难以被简单复制。", size=12.5, color=LGOLD, bold=True)

# ============ SLIDE 11: COMPETITION (table) ============
s = prs.slides.add_slide(BLANK)
add_bg(s, os.path.join(ASSET, "content-bg.png"))
badge(s, "竞品分析")
sub(s, "首个把『自动假设+证据可信+人在回路+闭环自迭代』统一的系统")
page_num(s, 11)
dims = ["科研理解", "自动假设生成", "证据可追溯性", "人工协同", "闭环自迭代"]
rows = [
    ("本项目", "强", "强", "强(白名单)", "强", "强"),
    ("ChatGPT", "弱", "无", "无", "无", "无"),
    ("Elicit", "中", "部分", "部分", "弱", "无"),
    ("AI Scientist", "中", "有", "弱", "弱", "部分"),
]
x0, w0, h0 = 0.5, 2.46, 0.85
for ci, head in enumerate(["维度", "本项目", "ChatGPT", "Elicit", "AI Scientist"]):
    x = x0 + ci*w0
    hd = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.5), Inches(w0-0.08), Inches(0.6))
    hd.fill.solid(); hd.fill.fore_color.rgb = NAVY if ci==0 else (EBLUE if ci==1 else RGBColor(0x3A,0x47,0x60))
    hd.line.fill.background()
    tf = hd.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = head; r.font.size = Pt(15); r.font.color.rgb = WHITE; r.font.bold = True; set_ea(r)
for ri, row in enumerate(rows):
    y = 2.2 + ri*h0
    for ci, val in enumerate(row):
        x = x0 + ci*w0
        cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w0-0.08), Inches(h0-0.06))
        is_us = (ri==0)
        is_dim = (ci==0)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x12,0x20,0x3C) if is_us else WHITE if is_dim else RGBColor(0x0F,0x1A,0x32)
        cell.line.color.rgb = CARDL; cell.line.width = Pt(1)
        tf = cell.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = val
        r.font.size = Pt(14); r.font.color.rgb = (CYAN if is_us else (NAVY if is_dim else LIGHT))
        r.font.bold = (is_us or is_dim); set_ea(r)
text_box(s, 0.5, 6.7, 12.3, 0.5, "竞品多在局部能力进展，AISci 是唯一全链路统一的科研自动化系统。", size=13, color=LGOLD, bold=True)

# ============ SLIDE 12: BUSINESS MODEL ============
s = prs.slides.add_slide(BLANK)
content(s, "商业模式", "三层变现，从低成本订阅到高客单私有化", "v12_pyramid.png", 12)
tiers = [
    ("高校版 SaaS", "1,999 ~ 9,999 元 / 年 / 账号", "低门槛、攒口碑、建标杆", GOLD),
    ("企业研发版", "10 万 ~ 100 万元 / 年", "高客单、撑营收、可复制", EBLUE),
    ("私有化部署", "项目制（部署+定制+运维）", "满足信创 / 数据不出域", CYAN),
]
for i, (t, price, note, c) in enumerate(tiers):
    y = 1.5 + i*1.7
    card(s, 0.4, y, 6.1, 1.5)
    text_box(s, 0.62, y+0.1, 3.6, 0.5, t, size=18, color=WHITE, bold=True)
    text_box(s, 4.2, y+0.1, 2.2, 0.5, price, size=15, color=c, bold=True, align=PP_ALIGN.RIGHT)
    gold_line(s, 0.62, y+0.62, 2.5)
    text_box(s, 0.62, y+0.72, 5.7, 0.6, note, size=13, color=LIGHT)
text_box(s, 0.4, 6.78, 6.2, 0.55, "收入逻辑：用户 → 产品 → 收入；已签合同金额最有说服力（请填真实数）。", size=12.5, color=SUB)

# ============ SLIDE 13: GTM ============
s = prs.slides.add_slide(BLANK)
content(s, "市场推广", "从科研生态切入，再走向产业", "v13_gtm.png", 13)
phases = [
    ("第一阶段 · 高校实验室验证", "标杆案例 + 口碑裂变，沉淀早期用户与论文背书", EBLUE),
    ("第二阶段 · 科研机构推广", "院系 / 研究院批量授权，建立机构级合作", CYAN),
    ("第三阶段 · 企业研发市场", "医药 / 新材料 / AI 行业 SaaS + 私有化", GOLD),
]
for i, (t, d, c) in enumerate(phases):
    y = 1.55 + i*1.7
    card(s, 0.4, y, 6.1, 1.5)
    text_box(s, 0.62, y+0.1, 5.7, 0.5, t, size=17, color=c, bold=True)
    gold_line(s, 0.62, y+0.62, 2.5)
    text_box(s, 0.62, y+0.72, 5.7, 0.65, d, size=13.5, color=LIGHT)
text_box(s, 0.4, 6.78, 6.2, 0.55, "GTM 关键词：标杆案例 · 学术 KOL · 信创渠道 · 生态合作。", size=12.5, color=SUB)

# ============ SLIDE 14: RESULTS ============
s = prs.slides.add_slide(BLANK)
content(s, "应用成果", "不是 Demo，是经测试验证的近产品系统", "v14_result.png", 14)
res = [
    ("Science 125 典型问题", "假设新颖性评分 ≥ 6.0，证据链 ≥ 1 轮迭代，References 无虚构引用", "数据来源：团队测试"),
    ("评分表系统", "6 组测试均生成完整评分表，总评分 88–118 分，维度 41–56 项", "数据来源：团队测试"),
    ("联邦学习仿真", "支持 local_pack / Flower / FedML 三种后端，脚本可执行", "数据来源：团队测试"),
]
for i, (t, d, src) in enumerate(res):
    y = 1.45 + i*1.5
    card(s, 0.4, y, 6.1, 1.35)
    text_box(s, 0.62, y+0.08, 5.7, 0.45, t, size=16, color=WHITE, bold=True)
    text_box(s, 0.62, y+0.52, 5.7, 0.7, d, size=13, color=LIGHT)
    text_box(s, 0.62, y+1.05, 5.7, 0.28, src, size=11, color=GOLD, bold=True)
text_box(s, 0.4, 6.25, 6.2, 1.0, "里程碑：多 Agent 系统 · 七阶段 Pipeline · 70+ Skill · Web 平台 · 国产适配 · 完整 Demo", size=12.5, color=LGOLD, bold=True, line_spacing=1.3)

# ============ SLIDE 15: TEAM ============
s = prs.slides.add_slide(BLANK)
content(s, "团队介绍", "AI 算法 × 科研领域 × 工程落地的复合战队", "v15_team.png", 15)
text_box(s, 0.4, 1.4, 6.1, 0.5, "团队顾问", size=17, color=WHITE, bold=True)
card(s, 0.4, 1.9, 6.1, 0.95)
text_box(s, 0.62, 2.0, 5.7, 0.8, "[姓名 / 头衔] —— 提供技术与资源支持（请填）", size=13.5, color=LIGHT, anchor=MSO_ANCHOR.MIDDLE)
text_box(s, 0.4, 3.0, 6.1, 0.5, "项目负责人", size=17, color=WHITE, bold=True)
card(s, 0.4, 3.5, 6.1, 0.95)
text_box(s, 0.62, 3.6, 5.7, 0.8, "[姓名] —— 技术负责人 / 大创项目主持人（请填国家级经历）", size=13.5, color=LIGHT, anchor=MSO_ANCHOR.MIDDLE)
text_box(s, 0.4, 4.6, 6.1, 0.5, "核心成员", size=17, color=WHITE, bold=True)
card(s, 0.4, 5.1, 6.1, 1.4)
multi_text(s, 0.62, 5.2, 5.7, 1.2, [
    "▸ [姓名] 主要负责 [具体分工]（请填竞赛获奖 + 分工）",
    "▸ [姓名] 主要负责 [具体分工]",
    "▸ [姓名] 主要负责 [具体分工]",
], size=12.5, color=LIGHT, line_spacing=1.25, space=4)
text_box(s, 0.4, 6.6, 6.2, 0.5, "在导师指导下完成系统从 0 到 1，已具备完整前后端与可复现性保障。", size=12.5, color=LGOLD, bold=True)

# ============ SLIDE 16: EDUCATION & SOCIAL ============
s = prs.slides.add_slide(BLANK)
content(s, "教育与社会价值", "以项目为牵引，四维反哺科研育人生态", "v16_education.png", 16)
dims4 = [
    ("课程融合", "AI for Science 案例进课堂，反哺课程建设", EBLUE),
    ("平台共建", "开放科研智能体平台，支撑教学实验", CYAN),
    ("科研育苗", "以项目培养本科生科研与工程能力", GOLD),
    ("标准反哺", "沉淀科研可审计流程，贡献方法论", LGOLD),
]
for i, (t, d, c) in enumerate(dims4):
    y = 1.5 + i*1.28
    card(s, 0.4, y, 6.1, 1.12)
    text_box(s, 0.62, y+0.1, 2.0, 0.9, t, size=16, color=c, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, 2.6, y+0.1, 3.8, 0.9, d, size=13, color=LIGHT, anchor=MSO_ANCHOR.MIDDLE)
text_box(s, 0.4, 6.7, 6.2, 0.6, "破『信息孤岛』，立『可信科研』；技术革新→产业升级→人才赋能→行业领航。", size=12.5, color=LGOLD, bold=True)

# ============ SLIDE 17: FINANCE ============
s = prs.slides.add_slide(BLANK)
content(s, "财务与融资", "锚定基石收入，迈向产业化", "v17_finance.png", 17)
text_box(s, 0.4, 1.4, 6.1, 0.5, "资金用途（请填真实比例）", size=16, color=WHITE, bold=True)
gold_line(s, 0.4, 1.92, 3.0)
uses = ["模型优化与推理成本", "科研数据与语料建设", "市场推广与品牌建设", "产品商业化与交付"]
for i, u in enumerate(uses):
    y = 2.1 + i*0.7
    card(s, 0.4, y, 6.1, 0.6)
    text_box(s, 0.6, y+0.02, 4.5, 0.55, f"▸ {u}", size=13.5, color=LIGHT, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, 5.0, y+0.02, 1.4, 0.55, "[  %]", size=13.5, color=GOLD, bold=True, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
text_box(s, 0.4, 5.2, 6.1, 0.5, "融资需求与里程碑（请填轮次 + 金额 + 目标）", size=16, color=WHITE, bold=True)
text_box(s, 0.4, 5.75, 6.1, 1.0, "本轮融资加速从『优秀科研系统』向『具产业化潜力的 AI 创业项目』跃迁。", size=13, color=LIGHT, line_spacing=1.35)
text_box(s, 0.4, 6.8, 6.2, 0.5, "⚠ 所有数字请本人核定，禁止 AI 编造。", size=12, color=RED, bold=True)

# ============ SLIDE 18: ROADMAP ============
s = prs.slides.add_slide(BLANK)
content(s, "未来规划", "三阶段，从高校试点到企业商业化", "v18_roadmap.png", 18)
phases = [
    ("Year 1", "高校试点验证", "沉淀标杆案例与口碑", EBLUE),
    ("Year 2", "科研机构合作", "拓展 Institution 产品线", CYAN),
    ("Year 3", "企业商业化", "SaaS + 私有化规模营收", GOLD),
]
for i, (y, t, d, c) in enumerate(phases):
    yy = 1.55 + i*1.7
    card(s, 0.4, yy, 6.1, 1.5)
    text_box(s, 0.62, yy+0.1, 2.2, 0.5, y, size=20, color=c, bold=True)
    text_box(s, 2.8, yy+0.1, 3.6, 0.5, t, size=16, color=WHITE, bold=True)
    gold_line(s, 0.62, yy+0.62, 2.5)
    text_box(s, 0.62, yy+0.72, 5.7, 0.6, d, size=13.5, color=LIGHT)
text_box(s, 0.4, 6.78, 6.2, 0.55, "目标用数字量化（请填营收 / 客户 / 案例等具体指标）。", size=12.5, color=SUB)

# ============ SLIDE 19: ENDING ============
s = prs.slides.add_slide(BLANK)
add_bg(s, os.path.join(ASSET, "thankyou-bg.png"))
gold_line(s, 3.6, 2.2, 6.1)
text_box(s, 0, 2.5, 13.33, 1.0, "联邦智研 AISci 之智，驱动科研创新新篇", size=38, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
text_box(s, 0, 3.7, 13.33, 0.7, "可信 · 可控 · 可复现 —— 让每个团队拥有 AI 科研协作者", size=22, color=LGOLD, bold=True, align=PP_ALIGN.CENTER)
gold_line(s, 4.1, 4.55, 5.1)
text_box(s, 0, 4.9, 13.33, 0.6, "感谢聆听，期待与您共筑 AI for Science 强国新章", size=20, color=WHITE, align=PP_ALIGN.CENTER)
text_box(s, 0, 6.3, 13.33, 0.5, "联邦智研团队  ·  [联系方式占位]", size=15, color=SUB, align=PP_ALIGN.CENTER)

prs.save(OUT)
print("SAVED", OUT, "slides:", len(prs.slides._sldIdLst))
