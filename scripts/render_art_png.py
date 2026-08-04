# -*- coding: utf-8 -*-
"""
把 v7.pptx 中的 OOXML 艺术字（gradFill/glow）改为【真实 PNG 艺术字图片】，
确保任何预览器 / PowerPoint / WPS 都能看到渐变+发光+阴影效果。
生成 output/AISci_互联网+答辩PPT_v8.pptx
"""
import os, shutil, json
from PIL import Image, ImageDraw, ImageFont, ImageFilter
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu

DPI = 200
PAD_IN = 0.28  # 图片四周留白（容纳发光/阴影）

FONT_PATH = r"C:/Windows/Fonts/HarmonyOS_Sans_SC_Bold.ttf"

# 配色（RGB）
CYAN_TOP = (0, 229, 255)      # #00E5FF
CYAN_BOT = (59, 130, 246)     # #3B82F6
GOLD_TOP = (255, 224, 102)    # #FFE066
GOLD_BOT = (255, 179, 0)      # #FFB300
GLOW_CYAN = (0, 229, 255)
GLOW_GOLD = (255, 196, 0)

A = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

def make_gradient(w, h, top, bottom):
    arr = np.zeros((h, w, 3), dtype=np.uint8)
    for c in range(3):
        arr[:, :, c] = np.linspace(top[c], bottom[c], h, dtype=np.uint8)[:, None]
    return Image.fromarray(arr, 'RGB')

def render_glyph(text, size_pt, colors, glow_color, align, box_w_in, box_h_in):
    """渲染单段文字为带渐变+发光+阴影的透明 PNG，返回 (PIL.Image, place_left_in, place_top_in)"""
    font_px = int(round(size_pt * DPI / 72.0))
    font = ImageFont.truetype(FONT_PATH, font_px)

    # 自适应字号：若超出文本框宽度则缩小
    tmp = Image.new('RGBA', (10, 10))
    d0 = ImageDraw.Draw(tmp)
    while True:
        bbox = d0.textbbox((0, 0), text, font=font)
        tw = bbox[2] - bbox[0]
        box_w_px = box_w_in * DPI
        if tw <= box_w_px * 0.97 or font_px <= 12:
            break
        font_px = int(font_px * 0.95)
        font = ImageFont.truetype(FONT_PATH, font_px)

    bbox = d0.textbbox((0, 0), text, font=font)
    tw = bbox[2] - bbox[0]
    th = bbox[3] - bbox[1]

    pad_px = int(PAD_IN * DPI)
    box_w_px = int(box_w_in * DPI)
    box_h_px = int(box_h_in * DPI)
    img_w = box_w_px + 2 * pad_px
    img_h = box_h_px + 2 * pad_px

    # 文字位置（图片内）
    if align == 'center':
        tx = (img_w - tw) / 2 - bbox[0]
    else:
        tx = pad_px - bbox[0]  # 左对齐贴文本框左边
    ty = (img_h - th) / 2 - bbox[1]

    # 1) 文字 alpha mask
    txt = Image.new('RGBA', (img_w, img_h), (0, 0, 0, 0))
    ImageDraw.Draw(txt).text((tx, ty), text, font=font, fill=(255, 255, 255, 255))
    mask = txt.split()[3]

    # 2) 渐变填充文字
    grad = make_gradient(img_w, img_h, colors[0], colors[1])
    colored = Image.new('RGBA', (img_w, img_h), (0, 0, 0, 0))
    colored.paste(grad, (0, 0))
    colored.putalpha(mask)

    # 3) 外发光（两层：亮晕 + 柔晕）
    glow = Image.new('RGBA', (img_w, img_h), (0, 0, 0, 0))
    g1 = mask.filter(ImageFilter.GaussianBlur(radius=max(4, font_px * 0.12)))
    glow.paste(glow_color + (255,), (0, 0), g1)
    g2 = mask.filter(ImageFilter.GaussianBlur(radius=max(10, font_px * 0.30)))
    glow2 = Image.new('RGBA', (img_w, img_h), (0, 0, 0, 0))
    glow2.paste(glow_color + (140,), (0, 0), g2)
    glow = Image.alpha_composite(glow2, glow)

    # 4) 阴影（偏移 + 模糊）
    shadow = Image.new('RGBA', (img_w, img_h), (0, 0, 0, 0))
    sm = mask.filter(ImageFilter.GaussianBlur(radius=max(5, font_px * 0.16)))
    sh_img = Image.new('RGBA', (img_w, img_h), (0, 0, 0, 0))
    sh_img.paste((0, 0, 0, 150), (0, 0), sm)
    shadow = sh_img.transform((img_w, img_h), Image.AFFINE, (1, 0, 0, 0, 1, max(3, font_px * 0.06)))

    # 合成
    base = Image.new('RGBA', (img_w, img_h), (0, 0, 0, 0))
    base = Image.alpha_composite(base, shadow)
    base = Image.alpha_composite(base, glow)
    base = Image.alpha_composite(base, colored)
    return base

def main():
    src = r"D:/Workplace/AISci/output/AISci_互联网+答辩PPT_v7.pptx"
    out = r"D:/Workplace/AISci/output/AISci_互联网+答辩PPT_v8.pptx"
    art_dir = r"D:/Workplace/AISci/output/_art_text"
    os.makedirs(art_dir, exist_ok=True)
    shutil.copy2(src, out)

    prs = Presentation(out)

    rendered = 0
    counter = 0
    for i, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            tf = sh.text_frame
            # 找艺术字 run
            art_runs = []
            for para in tf.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    rpr = run._r.find('.//' + A + 'rPr')
                    if rpr is None:
                        continue
                    if rpr.find(A + 'gradFill') is not None:
                        sz = round((run.font.size or 0) / 12700)
                        art_runs.append((run, sz))
            if not art_runs:
                continue

            counter += 1

            # 配色 & 对齐
            if i == 1 and art_runs[0][1] == 26:
                colors = (GOLD_TOP, GOLD_BOT); glow = GLOW_GOLD
            else:
                colors = (CYAN_TOP, CYAN_BOT); glow = GLOW_CYAN
            align = 'center' if i == 19 else 'left'

            left_in = (sh.left or 0) / 914400
            top_in = (sh.top or 0) / 914400
            width_in = (sh.width or 0) / 914400
            height_in = (sh.height or 0) / 914400

            # 段落文本拼接（保留原段落顺序，单行标题）
            full_text = tf.text.strip().replace('\n', ' ')

            img = render_glyph(full_text, art_runs[0][1], colors, glow, align,
                               width_in, height_in)
            fn = os.path.join(art_dir, f"art_{counter:03d}.png")
            img.save(fn, optimize=True)

            # 插入图片（覆盖原文本框位置 + padding）
            place_left = left_in - PAD_IN
            place_top = top_in - PAD_IN
            slide.shapes.add_picture(fn, Inches(place_left), Inches(place_top),
                                     Inches(img.width / DPI), Inches(img.height / DPI))

            # 清空原文本框文字（避免预览 fallback 黑色文字叠加）
            for para in tf.paragraphs:
                for run in para.runs:
                    if run._r.find('.//' + A + 'gradFill') is not None:
                        run.text = ''

            rendered += 1
            print(f"  S{i:02d} rendered art-text: '{full_text[:24]}' {art_runs[0][1]}pt align={align} -> {fn}")

    prs.save(out)
    print(f"\n✅ 完成！渲染 {rendered} 个艺术字标题 → {out}")
    print(f"   大小: {os.path.getsize(out)/1e6:.1f} MB")

if __name__ == '__main__':
    main()
