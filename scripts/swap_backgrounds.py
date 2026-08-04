#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Swap the plain PIL backgrounds of AISci_互联网+答辩PPT.pptx with the new
watermark-free 即梦 AI backgrounds (output/图片/第一页..第十九页), mapped 1:1 by
slide index. Downscale to 2000px wide to control file size, and add a subtle
dark scrim above the background (below content) for guaranteed text readability.
"""
import os, shutil
from io import BytesIO
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

PPTX = r"D:/Workplace/AISci/output/AISci_互联网+答辩PPT.pptx"
IMG_BASE = r"D:/Workplace/AISci/output/图片"
BACKUP = r"D:/Workplace/AISci/output/AISci_互联网+答辩PPT_备份.pptx"
MAXW = 2000
SCRIM_ALPHA = 18000  # 18% opaque black scrim (82% transparent) -> mild darkening

ZH = ['第一页','第二页','第三页','第四页','第五页','第六页','第七页','第八页','第九页',
      '第十页','第十一页','第十二页','第十三页','第十四页','第十五页','第十六页','第十七页','第十八页','第十九页']

def get_img(idx):
    d = os.path.join(IMG_BASE, ZH[idx-1])
    pngs = [f for f in os.listdir(d) if f.lower().endswith('.png')]
    if not pngs:
        return None
    p = os.path.join(d, pngs[0])
    im = Image.open(p).convert('RGB')
    if im.width > MAXW:
        h = int(im.height * MAXW / im.width)
        im = im.resize((MAXW, h), Image.LANCZOS)
    buf = BytesIO()
    im.save(buf, format='PNG')
    buf.seek(0)
    return buf

# backup original
if not os.path.exists(BACKUP):
    shutil.copy2(PPTX, BACKUP)
    print('backup ->', os.path.basename(BACKUP))

prs = Presentation(PPTX)
W = prs.slide_width
H = prs.slide_height

for idx, slide in enumerate(prs.slides, start=1):
    buf = get_img(idx)
    if buf is None:
        print(f'slide {idx}: NO IMAGE, skip')
        continue
    # locate the full-bleed background picture (left ~0, width ~full)
    bg_shape = None
    for sh in slide.shapes:
        if sh.shape_type == 13 and sh.left is not None:
            if abs(sh.left) < 10000 and sh.width and sh.width > W * 0.9:
                bg_shape = sh
                break
    if bg_shape is None:
        print(f'slide {idx}: no bg picture found, skip')
        continue
    # remove old bg
    sp = bg_shape._element
    sp.getparent().remove(sp)
    # add new bg at full bleed, then send to back
    pic = slide.shapes.add_picture(buf, 0, 0, W, H)
    psp = pic._element
    psp.getparent().remove(psp)
    slide.shapes._spTree.insert(2, psp)
    # subtle scrim above bg, below all content
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
    rect.fill.fore_color.alpha = SCRIM_ALPHA
    rect.line.fill.background()
    rsp = rect._element
    rsp.getparent().remove(rsp)
    slide.shapes._spTree.insert(3, rsp)
    print(f'slide {idx:2d}: bg swapped + scrim')

prs.save(PPTX)
print('SAVED', PPTX)
