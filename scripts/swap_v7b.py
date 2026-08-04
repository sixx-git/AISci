"""
v7b: 从 v4 模板替换背景图为 v6 修复后的干净版本。
用 python-pptx 标准 API 替换图片，避免缓存。
"""
import os, io, time
from PIL import Image
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

TEMPLATE = r"D:\Workplace\AISci\output\AISci_互联网+答辩PPT_v4.pptx"
OUTPUT   = r"D:\Workplace\AISci\output\AISci_互联网+答辩PPT_v5.pptx"

PAGE_MAP = {
    1:"第一页",2:"第二页",3:"第三页",4:"第四页",5:"第五页",
    6:"第六页",7:"第七页",8:"第八页",9:"第九页",10:"第十页",
    11:"第十一页",12:"第十二页",13:"第十三页",14:"第十四页",
    15:"第十五页",16:"第十六页",17:"第十七页",18:"第十八页",
    19:"第十九页",
}
IMG_BASE = r"D:\Workplace\AISci\output\图片"
MAX_W = 2000


def find_png(folder):
    d = os.path.join(IMG_BASE, folder)
    for f in os.listdir(d):
        if f.lower().endswith('.png'):
            return os.path.join(d, f)
    return None


def main():
    t0 = time.time()
    # 复制模板为新文件
    import shutil
    shutil.copy2(TEMPLATE, OUTPUT)

    prs = Presentation(OUTPUT)
    NS_B = '{http://schemas.openxmlformats.org/drawingml/2006/main}blip'
    NS_R = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'

    ok = 0
    for slide_idx, folder in PAGE_MAP.items():
        slide = prs.slides[slide_idx - 1]
        png_path = find_png(folder)
        if not png_path:
            continue

        # 强制从磁盘读原始字节（绕过 PIL 缓存）
        with open(png_path, 'rb') as f:
            raw_bytes = f.read()
        im = Image.open(io.BytesIO(raw_bytes)).convert('RGB')
        w, h = im.size
        if w > MAX_W:
            r = MAX_W / w
            im = im.resize((MAX_W, int(h * r)), Image.LANCZOS)

        buf = io.BytesIO()
        im.save(buf, format='PNG')
        img_data = buf.getvalue()

        # 找全出血背景 shape
        bg_sh = None
        for sh in slide.shapes:
            if sh.shape_type == 13 and sh.left is not None:
                wid = (sh.width or 0) / 914400
                if wid > 12:
                    bg_sh = sh; break
        if not bg_sh:
            continue

        # 用新图片替换：先添加到包中，再更新关系
        blip = bg_sh._element.find('.//' + NS_B)
        rid = blip.get(NS_R)

        # 添加新图片到 slide part 的 package
        new_part = slide.part.package.add_image(img_data)  # 尝试标准方法
        try:
            new_rid = slide.part.relate_to(new_part, RT.IMAGE)
        except:
            # fallback: 直接操作 XML
            from pptx.opc.package import Part
            from pptx.opc.packuri import PackURI
            image_num = len([r for r in slide.part.rels.values() if 'image' in str(r.reltype)]) + 1
            uri = PackURI(f'/ppt/media/image{image_num}.png')
            new_part = Part(uri, 'image/png', img_data, slide.part.package)
            slide.part.package.parts[uri] = new_part
            new_rid = slide.part.relate_to(new_part, RT.IMAGE)

        blip.set(NS_R, new_rid)
        ok += 1
        print(f"  ✅ {slide_idx:2d}({folder}) {im.size[0]}x{im.size[1]}")

    prs.save(OUTPUT)
    print(f"\n完成! {ok}/19 替换 | {os.path.getsize(OUTPUT)/1e6:.1f}MB | {time.time()-t0:.1f}s")


if __name__ == '__main__':
    main()
