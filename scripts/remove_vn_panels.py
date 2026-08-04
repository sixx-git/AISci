"""
remove_vn_panels.py — 移除所有 vN 不透明右侧面板，让即梦全出血背景成为主视觉

操作：
1. 打开 pptx
2. 对每张幻灯片，识别右侧视觉图（shape_type=PICTURE 且 width < 12 inch，即非全出血背景）
3. 删除该形状（从 XML 树中移除）
4. 保存
5. 验证：每张幻灯片只剩全出血背景图
"""

import sys, os, copy
from pptx import Presentation
from pptx.util import Emu, Inches

PPTX = os.path.join(os.path.dirname(__file__), "..", "output", "AISci_互联网+答辩PPT.pptx")
FULL_BLEED_MIN_W = Emu(Inches(12))   # 全出血背景宽度 > 12 inch；vN 视觉图 < 12 inch

prs = Presentation(PPTX)
print(f"打开 {PPTX}")
print(f"幻灯片总数: {len(prs.slides)}")
print(f"画布尺寸: {prs.slide_width/914400:.2f} x {prs.slide_height/914400:.2f} inch\n")

total_removed = 0
removed_log = []

for i, slide in enumerate(prs.slides, 1):
    to_remove = []
    for sh in slide.shapes:
        # 只处理图片形状，且排除全出血背景
        if sh.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            w = sh.width if sh.width is not None else Emu(0)
            if w < FULL_BLEED_MIN_W:
                to_remove.append(sh)

    if to_remove:
        for sh in to_remove:
            # 从父元素的 XML 中移除该形状
            sp = sh._element
            sp.getparent().remove(sp)
            total_removed += 1
        removed_log.append((i, len(to_remove)))
    else:
        removed_log.append((i, 0))

# 保存
prs.save(PPTX)
print(f"已保存。共删除 {total_removed} 个 vN 视觉图形状。\n")

# 验证：重新打开确认
prs2 = Presentation(PPTX)
print("=== 验证 ===")
all_clean = True
for i, slide in enumerate(prs2.slides, 1):
    pics = []
    for sh in slide.shapes:
        if sh.shape_type == 13 and sh.left is not None:
            w_inch = sh.width / 914400 if sh.width else 0
            kind = "BG" if w_inch > 12 else "VN!!"
            pics.append(f"{kind}({w_inch:.1f}\")")
    status = "OK" if all(p.startswith("BG") for p in pics) else "CHECK"
    if status == "CHECK":
        all_clean = False
    n_removed = removed_log[i-1][1]
    print(f"  Slide {i:2d}: 删除了 {n_removed} 个VN | 剩余图片={pics or '(无)'} [{status}]")

print(f"\n{'✅ 所有 VN 已清除' if all_clean else '⚠️ 有残留 VN'}")
print(f"\n文件大小: {os.path.getsize(PPTX)/1e6:.1f} MB")
