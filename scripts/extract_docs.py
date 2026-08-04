"""批量抽取 docx / pptx 正文到 txt，便于阅读设计文案。"""
import os
from docx import Document
from pptx import Presentation

OUT = r'D:\Workplace\AISci\output\_extract'
os.makedirs(OUT, exist_ok=True)


def dump_docx(path, name):
    d = Document(path)
    lines = []
    lines.append(f'===== {name} =====\n')
    # 段落
    for p in d.paragraphs:
        t = p.text.strip()
        if t:
            style = p.style.name if p.style else ''
            lines.append(f'[{style}] {t}')
    # 表格
    for ti, tb in enumerate(d.tables):
        lines.append(f'\n--- 表格 {ti+1} ({len(tb.rows)}x{len(tb.columns)}) ---')
        for row in tb.rows:
            cells = [c.text.strip().replace('\n', ' / ') for c in row.cells]
            lines.append(' | '.join(cells))
    txt = '\n'.join(lines)
    with open(os.path.join(OUT, name + '.txt'), 'w', encoding='utf-8') as f:
        f.write(txt)
    print(name, 'chars=', len(txt))


def dump_pptx(path, name):
    prs = Presentation(path)
    lines = [f'===== {name} (共 {len(prs.slides)} 页) =====\n']
    for i, sld in enumerate(prs.slides, 1):
        lines.append(f'\n########## 第 {i} 页 ##########')
        for sh in sld.shapes:
            st = sh.shape_type
            txt = ''
            if sh.has_text_frame:
                txt = sh.text_frame.text.strip()
            if txt:
                lines.append(f'[{st}] {txt}')
            if sh.has_table:
                for row in sh.table.rows:
                    cells = [c.text.strip().replace('\n', ' / ') for c in row.cells]
                    lines.append('   TBL | ' + ' | '.join(cells))
    txt = '\n'.join(lines)
    with open(os.path.join(OUT, name + '.txt'), 'w', encoding='utf-8') as f:
        f.write(txt)
    print(name, 'chars=', len(txt))


base = r'D:\Workplace\AISci\output'
dump_docx(os.path.join(base, '联邦智研_大创申报书_v4.docx'), 'shenbao_v4')
dump_docx(os.path.join(base, '项目计划书-联邦智研AISci-商业版.docx'), 'business_plan')
dump_pptx(os.path.join(base, 'AISci_互联网+答辩PPT_v12.pptx'), 'pptx_v12')
print('done')
