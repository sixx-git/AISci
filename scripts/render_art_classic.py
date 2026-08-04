# -*- coding: utf-8 -*-
"""
v9：经典 WordArt 风格艺术字（PNG）—— 去发光、加描边、柔和投影、柔化渐变。
解决 v8「晃眼」问题：真正的艺术字观感，且在任何预览器都可见。
基于 v7（带 gradFill 的文本框）复制 → 渲染经典 PNG → 插入 → 清空原文本。
输出 output/AISci_互联网+答辩PPT_v9.pptx
"""
import os, shutil
from PIL import Image, ImageDraw, ImageFont, ImageFilter
import numpy as np
from scipy.ndimage import binary_dilation
from pptx import Presentation
from pptx.util import Inches

DPI = 200
PAD_IN = 0.30

FONT_PATH = r"C:/Windows/Fonts/HarmonyOS_Sans_SC_Bold.ttf"

# 柔和渐变（比 v8 降饱和/亮度，避免刺眼）
CYAN_TOP = (43, 212, 240)     # #2BD4F0
CYAN_BOT = (47, 111, 224)     # #2F6FE0
GOLD_TOP = (255, 222, 133)    # #FFDE85
GOLD_BOT = (239, 166, 60)     # #EFA63C
# 描边（同色系深色，包裹亮色边缘，增强质感、消除刺眼亮边）
STROKE_CYAN = (10, 61, 107)   # #0A3D6B
STROKE_GOLD = (138, 90, 18)   # #8A5A12

A = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

def make_gradient(w, h, top, bottom):
    arr = np.zeros((h, w, 3), dtype=np.uint8)
    for c in range(3):
        arr[:, :, c] = np.linspace(top[c], bottom[c], h, dtype=np.uint8)[:, None]
    return Image.fromarray(arr, 'RGB')

def render_glyph(text, size_pt, colors, stroke_color, align, box_w_in, box_h_in):
    font_px = int(round(size_pt * DPI / 72.0))
    font = ImageFont.truetype(FONT_PATH, font_px)
    tmp = Image.new('RGBA', (10, 10)); d0 = ImageDraw.Draw(tmp)
    while True:
        bbox = d0.textbbox((0, 0), text, font=font)
        tw = bbox[2] - bbox[0]
        if tw <= box_w_in * DPI * 0.97 or font_px <= 12:
            break
        font_px = int(font_px * 0.95)
        font = ImageFont.truetype(FONT_PATH, font_px)

    bbox = d0.textbbox((0, 0), text, font=font)
    tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
    pad_px = int(PAD_IN * DPI)
    img_w = int(box_w_in * DPI) + 2 * pad_px
    img_h = int(box_h_in * DPI) + 2 * pad_px

    if align == 'center':
        tx = (img_w - tw) / 2 - bbox[0]
    else:
        tx = pad_px - bbox[0]
    ty = (img_h - th) / 2 - bbox[1]

    # 文字 alpha
    txt = Image.new('RGBA', (img_w, img_h), (0, 0, 0, 0))
    ImageDraw.Draw(txt).text((tx, ty), text, font=font, fill=(255, 255, 255, 255))
    mask = np.array(txt.split()[3])

    # 描边（膨胀 - 原mask）
    stroke_px = max(2, int(font_px * 0.028))
    dil = binary_dilation(mask > 0, iterations=stroke_px)
    stroke_np = dil & (~(mask > 0))
    stroke_mask = Image.fromarray((stroke_np * 255).astype('uint8'), 'L')
    stroke_layer = Image.new('RGBA', (img_w, img_h), (0, 0, 0, 0))
    stroke_layer.paste(stroke_color + (255,), (0, 0), stroke_mask)

    # 渐变填充文字
    grad = make_gradient(img_w, img_h, colors[0], colors[1])
    colored = Image.new('RGBA', (img_w, img_h), (0, 0, 0, 0))
    colored.paste(grad, (0, 0))
    colored.putalpha(Image.fromarray(mask, 'L'))

    # 柔和投影（弱）
    mask_img = Image.fromarray(mask, 'L')
    shadow = Image.new('RGBA', (img_w, img_h), (0, 0, 0, 0))
    sm = mask_img.filter(ImageFilter.GaussianBlur(radius=max(4, font_px * 0.07)))
    sh_img = Image.new('RGBA', (img_w, img_h), (0, 0, 0, 0))
    sh_img.paste((0, 0, 0, 110), (0, 0), sm)
    shadow = sh_img.transform((img_w, img_h), Image.AFFINE,
                              (1, 0, 0, 0, 1, max(2, font_px * 0.03)))

    base = Image.new('RGBA', (img_w, img_h), (0, 0, 0, 0))
    base = Image.alpha_composite(base, shadow)
    base = Image.alpha_composite(base, stroke_layer)
    base = Image.alpha_composite(base, colored)
    return base

def main():
    src = r"D:/Workplace/AISci/output/AISci_互联网+答辩PPT_v7.pptx"
    out = r"D:/Workplace/AISci/output/AISci_互联网+答辩PPT_v9.pptx"
    art_dir = r"D:/Workplace/AISci/output/_art_text_v9"
    os.makedirs(art_dir, exist_ok=True)
    shutil.copy2(src, out)
    prs = Presentation(out)

    rendered = 0; counter = 0
    for i, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            tf = sh.text_frame
            art_runs = []
            for para in tf.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    rpr = run._r.find('.//' + A + 'rPr')
                    if rpr is None:
                        continue
                    if rpr.find(A + 'gradFill') is not None:
                        art_runs.append((run, round((run.font.size or 0) / 12700)))
            if not art_runs:
                continue
            counter += 1

            if i == 1 and art_runs[0][1] == 26:
                colors, stroke_color = (GOLD_TOP, GOLD_BOT), STROKE_GOLD
            else:
                colors, stroke_color = (CYAN_TOP, CYAN_BOT), STROKE_CYAN
            align = 'center' if i == 19 else 'left'

            left_in = (sh.left or 0) / 914400
            top_in = (sh.top or 0) / 914400
            width_in = (sh.width or 0) / 914400
            height_in = (sh.height or 0) / 914400
            full_text = tf.text.strip().replace('\n', ' ')

            img = render_glyph(full_text, art_runs[0][1], colors, stroke_color,
                               align, width_in, height_in)
            fn = os.path.join(art_dir, f"art_{counter:03d}.png")
            img.save(fn, optimize=True)
            slide.shapes.add_picture(fn, Inches(left_in - PAD_IN), Inches(top_in - PAD_IN),
                                     Inches(img.width / DPI), Inches(img.height / DPI))
            for para in tf.paragraphs:
                for run in para.runs:
                    if run._r.find('.//' + A + 'gradFill') is not None:
                        run.text = ''
            rendered += 1
            print(f"  S{i:02d} [{align}] '{full_text[:24]}' {art_runs[0][1]}pt -> {fn}")

    prs.save(out)
    print(f"\n✅ 完成！{rendered} 个经典艺术字 → {out} ({os.path.getsize(out)/1e6:.1f} MB)")

if __name__ == '__main__':
    main()
