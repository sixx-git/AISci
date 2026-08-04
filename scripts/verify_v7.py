"""Verify v7 PPT: backgrounds, art-text XML, accent bars."""
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

A = 'http://schemas.openxmlformats.org/drawingml/2006/main'

p = r'D:\Workplace\AISci\output\AISci_互联网+答辩PPT_v7.pptx'
prs = Presentation(p)
d = r'D:\Workplace\AISci\_cmp\v7_verify'
os.makedirs(d, exist_ok=True)

# 1) Extract backgrounds for key slides
for sidx in [1, 3, 10, 19]:
    slide = prs.slides[sidx - 1]
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE and sh.left is not None and abs(sh.left) < 10000:
            w = (sh.width or 0) / 914400
            if w > 12:
                outp = os.path.join(d, f'bg_{sidx:02d}.{sh.image.ext}')
                with open(outp, 'wb') as f:
                    f.write(sh.image.blob)
                print(f'BG p{sidx:02d}: {outp} ({len(sh.image.blob)//1024}KB)')
                break

# 2) Verify art-text XML on titles
print('\n--- Art-text verification ---')
checks = [
    (1, '联邦智研 AISci', 'cover main'),
    (3, '政策背景', 'section'),
    (10, '技术壁垒', 'section'),
    (19, '联邦智研 AISci 之智', 'ending'),
]
for sidx, txt_frag, label in checks:
    slide = prs.slides[sidx - 1]
    found = False
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                if txt_frag in run.text:
                    rpr = run._r.find('.//{%s}rPr' % A)
                    if rpr is not None:
                        gf = rpr.find('{%s}gradFill' % A)
                        eff = rpr.find('{%s}effectLst' % A)
                        gl = eff.find('{%s}glow' % A) if eff is not None else None
                        shdw = eff.find('{%s}outerShdw' % A) if eff is not None else None
                        print(f'  p{sidx:02d} [{label}] "{run.text[:30]}" '
                              f'gradFill={"YES" if gf else "NO"} '
                              f'glow={"YES" if gl else "NO"} '
                              f'shadow={"YES" if shdw else "NO"}')
                        found = True
    if not found:
        print(f'  p{sidx:02d} [{label}] NOT FOUND: "{txt_frag}"')

# 3) Count accent bars
bars = 0
for i, slide in enumerate(prs.slides, 1):
    for sh in slide.shapes:
        if sh.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        w = (sh.width or 0) / 914400
        h = (sh.height or 0) / 914400
        left = (sh.left or 0) / 914400
        if 0.04 < w < 0.12 and 0.5 < h < 0.8 and abs(left - 0.05) < 0.03:
            bars += 1
print(f'\nAccent bars found: {bars}/17')
print('\nDone.')
