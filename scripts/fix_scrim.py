"""修复：将盖在背景图上的不透明黑色 scrim 改为 40% 半透明，使即梦背景透出。
同时强制背景图全出血铺满整页。"""
import os
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

SRC = "output/AISci_互联网+答辩PPT_clean.pptx"
OUT = "output/AISci_互联网+答辩PPT_v4.pptx"

prs = Presentation(SRC)
SW, SH = prs.slide_width, prs.slide_height
W_IN, H_IN = SW/914400, SH/914400
print(f"画布: {W_IN:.2f} x {H_IN:.2f} inch")

NS_B = '{http://schemas.openxmlformats.org/drawingml/2006/main}blip'
NS_R = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'

fixed_scrim = 0
fixed_bg = 0
for sidx, slide in enumerate(prs.slides, 1):
    # 1) 确保背景图全出血铺满
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE and sh.left is not None and abs(sh.left) < 10000:
            wid = sh.width/914400 if sh.width else 0
            if wid > 12:
                if sh.left != 0 or sh.top != 0 or abs(sh.width-SW) > 5000 or abs(sh.height-SH) > 5000:
                    sh.left = 0; sh.top = 0; sh.width = SW; sh.height = SH
                    fixed_bg += 1
    # 2) 找到全出血的黑色 scrim 矩形，加入 40% alpha
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            wid = sh.width/914400 if sh.width else 0
            if wid <= 12:
                continue  # 只处理全出血矩形
            # 找黑色 srgbClr
            srgb = sh._element.find(f'.//{qn("a:solidFill")}/{qn("a:srgbClr")}')
            if srgb is None:
                # 可能没有 solidFill，跳过
                continue
            if srgb.get('val') != '000000':
                continue
            # 已有 alpha 则改之，否则新增
            alpha = srgb.find(qn('a:alpha'))
            if alpha is None:
                alpha = srgb.makeelement(qn('a:alpha'), {})
                srgb.append(alpha)
            alpha.set('val', '40000')  # 40% 不透明 -> 图像透出 60%
            fixed_scrim += 1

prs.save(OUT)
print(f"✅ 修复完成: 调整背景尺寸 {fixed_bg} 处, 添加/修正 scrim 透明度 {fixed_scrim} 处")
print(f"输出: {OUT}  ({os.path.getsize(OUT)/1e6:.1f} MB)")
