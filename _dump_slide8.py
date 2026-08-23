from pptx import Presentation
from pptx.shapes.group import GroupShape

def dump_shape(shp, indent=0):
    sp = '  '*indent
    print(f"{sp}name: {shp.name}, type: {shp.shape_type}")
    print(f"{sp}  pos: ({shp.left.inches:.2f}, {shp.top.inches:.2f}) size: ({shp.width.inches:.2f}, {shp.height.inches:.2f})")
    if shp.has_text_frame:
        for j, para in enumerate(shp.text_frame.paragraphs):
            txt = ''.join([r.text for r in para.runs])
            if txt.strip():
                print(f"{sp}  txt: {txt[:140]}")
    if isinstance(shp, GroupShape):
        for child in shp.shapes:
            dump_shape(child, indent+2)

p = r"D:/Workplace/AISci/output/蓝白色医疗器械--淘宝店：烦烦优创馆.pptx"
prs = Presentation(p)
print(f"Total slides: {len(prs.slides)}")
print(f"\n=== Slide 8 (index 7) ===")
for i, shp in enumerate(prs.slides[7].shapes):
    print(f"\n--- top shape {i} ---")
    dump_shape(shp)
