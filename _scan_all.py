from pptx import Presentation

p = r"D:/Workplace/AISci/output/蓝白色医疗器械--淘宝店：烦烦优创馆.pptx"
prs = Presentation(p)

for si, slide in enumerate(prs.slides):
    texts = []
    npic = 0
    nchart = 0
    for shp in slide.shapes:
        if shp.shape_type == 13:
            npic += 1
        if shp.shape_type == 3:
            nchart += 1
        if shp.has_text_frame:
            for para in shp.text_frame.paragraphs:
                t = ''.join([r.text for r in para.runs]).strip()
                if t and len(t) < 60:
                    texts.append(t)
    # dedupe keep first 6
    seen = []
    for t in texts:
        if t not in seen:
            seen.append(t)
    print(f"\n===== Slide {si+1} (pics={npic}, charts={nchart}) =====")
    for t in seen[:8]:
        print(f"  · {t}")
