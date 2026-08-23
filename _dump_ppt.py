import docx
from pptx import Presentation
from pptx.util import Emu

p = r"D:/Workplace/AISci/output/蓝白色医疗器械--淘宝店：烦烦优创馆.pptx"
prs = Presentation(p)

print("=== 幻灯片尺寸 ===")
print("slide_width:", Emu(prs.slide_width).inches, "in x", Emu(prs.slide_height).inches, "in")

for idx, slide in enumerate(prs.slides, 1):
    print(f"\n{'='*70}")
    print(f"### 第 {idx} 页 ###")
    layout = slide.slide_layout
    print(f"版式(layout): {layout.name}")
    # background
    bg = slide.background
    try:
        if bg.fill.type is not None:
            print(f"背景填充类型: {bg.fill.type}")
    except Exception:
        pass
    n_shapes = len(slide.shapes)
    print(f"形状数量: {n_shapes}")
    for si, shape in enumerate(slide.shapes):
        stype = shape.shape_type
        has_tf = shape.has_text_frame
        txt = ""
        if has_tf:
            txt = shape.text_frame.text.strip().replace("\n", " / ")
        img = " [图片]" if shape.shape_type == 13 else ""
        print(f"  - shape[{si}] type={stype}{img} txt={txt[:160]}")
