# -*- coding: utf-8 -*-
"""Expanded PPT: put the report's text & tables into slides, presentation-friendly.
Chapters kept as groups; each chapter may span multiple slides.
"""
import os
from PIL import Image, ImageDraw, ImageFont, ImageFilter
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
def OxmlElement(tag):
    return etree.Element(qn(tag))

ASSET = r"D:\Workplace\AISci\_ppt_assets"
OUT_PPTX = r"D:\Workplace\AISci\output\联邦智研AISci_赛题指标达成与阶段进展报告_PPT.pptx"
os.makedirs(ASSET, exist_ok=True)
W, H = 2000, 1125

NAVY = RGBColor(0x0A, 0x16, 0x28)
DARK_BLUE = RGBColor(0x0F, 0x1D, 0x35)
GOLD = RGBColor(0xC9, 0xA8, 0x4C)
LIGHT_GOLD = RGBColor(0xE8, 0xD5, 0x8C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
SUB_TEXT = RGBColor(0x55, 0x66, 0x80)
ACCENT = RGBColor(0x38, 0x6F, 0xC4)
GREEN = RGBColor(0x1E, 0x7A, 0x46)
RED = RGBColor(0xC0, 0x39, 0x2B)
GRAY_BODY = RGBColor(0x4B, 0x55, 0x63)
LIGHT_BORDER = RGBColor(0xE2, 0xE8, 0xF0)
LIGHT_BORDER_B = RGBColor(0xD2, 0xE0, 0xF2)
GRAD_BODY = RGBColor(0xF4, 0xF8, 0xFD)
GRAD_BODY2 = RGBColor(0xEE, 0xF3, 0xFB)
TECH_BLUE = RGBColor(0x38, 0x9B, 0xFF)
TEAL = RGBColor(0x2D, 0xD4, 0xBF)
GOLD_LINE = RGBColor(0xC5, 0xA0, 0x59)
INK_BLUE = RGBColor(0x10, 0x2A, 0x43)
SUB_LIGHT = RGBColor(0x8D, 0xA7, 0xBE)
SERIF_FONT = "Noto Serif SC"
SANS_FONT = "Noto Sans SC"
HEAVY_SANS = "HarmonyOS Sans SC"
NUM_FONT = "Bahnschrift"
FONT = "等线"

# ---------- backgrounds (reuse) ----------
def vgradient(d, w, h, stops):
    for y in range(h):
        ratio = y / h
        for i in range(len(stops) - 1):
            p0, c0 = stops[i]; p1, c1 = stops[i + 1]
            if p0 <= ratio <= p1:
                t = (ratio - p0) / (p1 - p0) if p1 > p0 else 0
                r = int(c0[0] + (c1[0] - c0[0]) * t)
                g = int(c0[1] + (c1[1] - c0[1]) * t)
                b = int(c0[2] + (c1[2] - c0[2]) * t)
                d.line([(0, y), (w, y)], fill=(r, g, b))
                break

def make_cover():
    img = Image.new("RGB", (W, H)); d = ImageDraw.Draw(img)
    vgradient(d, W, H, [(0.0, (8, 16, 34)), (0.55, (15, 28, 56)), (1.0, (24, 42, 78))])
    glow = Image.new("RGBA", (W, H), (0, 0, 0, 0)); gd = ImageDraw.Draw(glow)
    cx, cy = int(W * 0.82), int(H * 0.22)
    for r in range(620, 0, -18):
        gd.ellipse([cx - r, cy - r, cx + r, cy + r], fill=(201, 168, 76, int(46 * (1 - r / 620))))
    cx2, cy2 = int(W * 0.12), int(H * 0.9)
    for r in range(520, 0, -18):
        gd.ellipse([cx2 - r, cy2 - r, cx2 + r, cy2 + r], fill=(56, 111, 196, int(34 * (1 - r / 520))))
    img = Image.alpha_composite(img.convert("RGBA"), glow).convert("RGB")
    nd = ImageDraw.Draw(img); import random; random.seed(7)
    nodes = [(random.randint(0, W), random.randint(0, H)) for _ in range(10)]
    for i in range(len(nodes)):
        for j in range(i + 1, len(nodes)):
            if (nodes[i][0] - nodes[j][0]) ** 2 + (nodes[i][1] - nodes[j][1]) ** 2 < 360000:
                nd.line([nodes[i], nodes[j]], fill=(120, 150, 200), width=1)
    for (x, y) in nodes:
        nd.ellipse([x - 4, y - 4, x + 4, y + 4], fill=(201, 168, 76))
    img.save(os.path.join(ASSET, "cover-bg.png"))

def make_content():
    img = Image.new("RGB", (W, H)); d = ImageDraw.Draw(img)
    vgradient(d, W, H, [(0.0, (246, 248, 252)), (1.0, (232, 237, 245))])
    grid = Image.new("RGBA", (W, H), (0, 0, 0, 0)); gd = ImageDraw.Draw(grid)
    for x in range(0, W, 50): gd.line([(x, 0), (x, H)], fill=(180, 195, 220, 22))
    for y in range(0, H, 50): gd.line([(0, y), (W, y)], fill=(180, 195, 220, 22))
    img = Image.alpha_composite(img.convert("RGBA"), grid).convert("RGB")
    glow = Image.new("RGBA", (W, H), (0, 0, 0, 0)); gd = ImageDraw.Draw(glow)
    cx, cy = int(W * 0.98), int(H * 0.98)
    for r in range(420, 0, -16):
        gd.ellipse([cx - r, cy - r, cx + r, cy + r], fill=(56, 111, 196, int(30 * (1 - r / 420))))
    img = Image.alpha_composite(img.convert("RGBA"), glow).convert("RGB")
    img.save(os.path.join(ASSET, "content-bg.png"))

make_cover(); make_content()
CB = os.path.join(ASSET, "cover-bg.png")
COB = os.path.join(ASSET, "content-bg.png")

# ---------- cinematic backgrounds sourced from v12 答辩PPT (cleaner, smaller mosaic) ----------
# 与 output/图片 同场景，但 v12 嵌入版右下角马赛克更小；按场景对应：
#   cover(第一页)→v12 slide1, 第十一页→v12 slide11, 第九页→v12 slide9, 第十九页→v12 slide19
V12_BASE = r"D:\Workplace\AISci\_ppt_assets\v12matched"
def load_bg(src_png, out_name):
    im = Image.open(src_png).convert("RGB").resize((2000, 1125), Image.LANCZOS)
    out = os.path.join(ASSET, out_name); im.save(out); return out
COVER_IMG = load_bg(os.path.join(V12_BASE, "bg_cover.png"), "bg_cover.png")
BG_P4 = load_bg(os.path.join(V12_BASE, "bg_p4.png"), "bg_p4.png")
BG_P5 = load_bg(os.path.join(V12_BASE, "bg_p5.png"), "bg_p5.png")
BG_P7 = load_bg(os.path.join(V12_BASE, "bg_p7.png"), "bg_p7.png")

# ---------- 水墨艺术字：书法标题渲染为透明 PNG ----------
def _load_brush_font(size):
    for fp in (r"C:\Windows\Fonts\STXINGKA.TTF",
               r"C:\Windows\Fonts\STKAITI.TTF",
               r"C:\Windows\Fonts\simkai.ttf"):
        if os.path.exists(fp):
            return ImageFont.truetype(fp, size)
    return ImageFont.load_default()

def ink_title(text, out_name, dark=False):
    fs = 150
    font = _load_brush_font(fs)
    tmp = Image.new("RGBA", (10, 10)); td = ImageDraw.Draw(tmp)
    bbox = td.textbbox((0, 0), text, font=font)
    tw = bbox[2] - bbox[0]; th = bbox[3] - bbox[1]
    pad_x, pad_y = 84, 34
    iw = int(tw + pad_x * 2); ih = int(th + pad_y * 2)
    img = Image.new("RGBA", (iw, ih), (0, 0, 0, 0))
    d = ImageDraw.Draw(img)
    ink = (238, 243, 250, 255) if dark else (20, 28, 44, 255)
    for (ox, oy, a) in [(-2, -2, 55), (2, 2, 55), (0, 3, 45), (-3, 1, 38)]:
        d.text((pad_x + ox - bbox[0], pad_y + oy - bbox[1]), text, font=font,
               fill=(ink[0], ink[1], ink[2], a))
    d.text((pad_x - bbox[0], pad_y - bbox[1]), text, font=font, fill=ink)
    img = img.filter(ImageFilter.GaussianBlur(0.7))
    seal = int(fs * 0.42)
    s = Image.new("RGBA", (seal, seal), (0, 0, 0, 0))
    sd = ImageDraw.Draw(s); red = (192, 57, 43, 240)
    sd.rectangle([3, 3, seal - 3, seal - 3], outline=red, width=4)
    sf = _load_brush_font(int(seal * 0.6))
    sb = sd.textbbox((0, 0), "研", font=sf)
    sd.text(((seal - (sb[2] - sb[0])) / 2 - sb[0], (seal - (sb[3] - sb[1])) / 2 - sb[1]),
            "研", font=sf, fill=red)
    img.paste(s, (iw - seal - 10, int((ih - seal) / 2)), s)
    out = os.path.join(ASSET, out_name); img.save(out)
    return out, iw, ih

_INK_CACHE = {}
def ink_badge(slide, text, dark=False):
    key = (text, dark)
    if key not in _INK_CACHE:
        _INK_CACHE[key] = ink_title(text, f"ink_{'d' if dark else 'l'}_{abs(hash(text))}.png", dark=dark)
    path, iw, ih = _INK_CACHE[key]
    target_h = 0.60
    w = target_h * (iw / ih)
    slide.shapes.add_picture(path, Inches(0.35), Inches(0.16), Inches(w), Inches(target_h))

# ---------- pptx helpers ----------
prs = Presentation()
prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = Inches(13.33), Inches(7.5)

def set_font(run, name=FONT):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", name)

def add_bg(slide, path):
    slide.background.fill.background()
    pic = slide.shapes.add_picture(path, Emu(0), Emu(0), SW, SH)
    sp = pic._element; sp.getparent().remove(sp); slide.shapes._spTree.insert(2, sp)

def text_box(slide, left, top, width, height, text, size=18, color=DARK_TEXT,
             bold=False, align=PP_ALIGN.LEFT, line_spacing=1.2, font=FONT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold; set_font(r, font)
    p.line_spacing = line_spacing
    return tb

def bullets(slide, left, top, width, height, items, size=14, color=DARK_TEXT, gap=6, head_color=NAVY, font=FONT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        head, body = it if isinstance(it, tuple) else (None, it)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT; p.space_after = Pt(gap); p.line_spacing = 1.12
        if head:
            r = p.add_run(); r.text = "● " + head; r.font.size = Pt(size); r.font.color.rgb = head_color; r.font.bold = True; set_font(r, font)
            r2 = p.add_run(); r2.text = "　" + body; r2.font.size = Pt(size); r2.font.color.rgb = color; set_font(r2, font)
        else:
            r = p.add_run(); r.text = "● " + body; r.font.size = Pt(size); r.font.color.rgb = color; set_font(r, font)
    return tb

def gold_line(slide, left, top, width, thick=2.5):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(thick))
    s.fill.solid(); s.fill.fore_color.rgb = GOLD; s.line.fill.background()

def badge(slide, text, left=0.35, top=0.28, w=2.7, h=0.58, size=19):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = NAVY; s.line.fill.background()
    tf = s.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.size = Pt(size); r.font.color.rgb = WHITE; r.font.bold = True; set_font(r)
    return s

def sub_header(slide, text, left=3.3, top=0.3, width=9.6, size=15, color=SUB_TEXT, font=FONT):
    text_box(slide, left, top, width, 0.6, "——" + text, size=size, color=color, font=font)

def page_num(slide, n, total, color=SUB_TEXT):
    text_box(slide, 12.1, 6.95, 1.0, 0.4, f"{n}/{total}", size=11, color=color, align=PP_ALIGN.RIGHT)

LIGHT_TEXT = RGBColor(0xE6, 0xEC, 0xF5)
LIGHT_SUB = RGBColor(0x9F, 0xB3, 0xD1)
def _ensure_effectLst(rPr):
    effectLst = rPr.find(qn('a:effectLst'))
    if effectLst is None:
        effectLst = OxmlElement('a:effectLst')
        idx = len(rPr)
        for i, ch in enumerate(rPr):
            t = ch.tag.split('}')[-1]
            if t in ('latin', 'ea', 'cs', 'highlight', 'uLn', 'uLnTx', 'uFill', 'uFillTx', 'sym'):
                idx = i; break
        rPr.insert(idx, effectLst)
    return effectLst

def _glow_run(run, rgb, alpha_pct, blur_pt):
    """文字外发光（极轻蓝发光）。"""
    rPr = run._r.get_or_add_rPr()
    effectLst = _ensure_effectLst(rPr)
    glow = OxmlElement('a:glow'); glow.set('rad', str(int(blur_pt * 12700)))
    c = OxmlElement('a:srgbClr'); c.set('val', '%02X%02X%02X' % rgb)
    a = OxmlElement('a:alpha'); a.set('val', str(int(alpha_pct * 1000))); c.append(a)
    glow.append(c); effectLst.append(glow)

def _outline_run(run, rgb, w_pt):
    """文字细描边（深蓝）。"""
    rPr = run._r.get_or_add_rPr()
    ln = OxmlElement('a:ln'); ln.set('w', str(int(w_pt * 12700)))
    sf = OxmlElement('a:solidFill'); c = OxmlElement('a:srgbClr'); c.set('val', '%02X%02X%02X' % rgb); sf.append(c)
    ln.append(sf); rPr.insert(0, ln)

def draw_title(slide, text, kind):
    """主标题渲染：
    academic = 思源宋体 Bold 深墨蓝 #102A43 + 1.5px 金色细线 #C5A059；
    darktech = 鸿蒙黑体 Bold 冰蓝白 #E6F0FF + 极轻蓝发光 + 深蓝细描边 + 底部蓝→金渐变线 + 左侧 3px 竖线。"""
    if kind == 'academic':
        font, color, size, bold = SERIF_FONT, INK_BLUE, 30, True
    else:
        font, color, size, bold = HEAVY_SANS, RGBColor(0xE6, 0xF0, 0xFF), 30, True
    tb = slide.shapes.add_textbox(Inches(0.35), Inches(0.2), Inches(5.4), Inches(0.72))
    tf = tb.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; set_font(r, font)
    if kind == 'academic':
        gl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.37), Inches(0.95), Inches(2.7), Pt(1.5))
        gl.fill.solid(); gl.fill.fore_color.rgb = GOLD_LINE; gl.line.fill.background()
    else:
        _glow_run(r, (38, 155, 255), 25, 12)
        _outline_run(r, (7, 16, 31), 0.75)
        gline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.37), Inches(0.95), Inches(2.9), Pt(1.5))
        _gradient_fill(gline, (56, 155, 255), (0xC5, 0xA0, 0x59), angle=0)
        gline.line.fill.background()
        ll = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.0), Inches(0.2), Pt(3), Inches(0.72))
        ll.fill.solid(); ll.fill.fore_color.rgb = TECH_BLUE; ll.line.fill.background()
    return tb

def content_slide(badge_text, subtitle, n, total, bg=COB, dark=False, title_kind="ink", sub_color=None, sub_font=None):
    s = prs.slides.add_slide(BLANK)
    add_bg(s, bg)
    if title_kind == "ink":
        ink_badge(s, badge_text, dark=dark)
    else:
        draw_title(s, badge_text, title_kind)
    sc = sub_color if sub_color is not None else (LIGHT_GOLD if dark else SUB_TEXT)
    sub_header(s, subtitle, color=sc, font=(sub_font or FONT))
    gold_line(s, 0.35, 1.0, 12.6, 2)
    page_num(s, n, total, color=(LIGHT_SUB if dark else SUB_TEXT))
    return s

def _gradient_fill(shape, top, bottom, angle=5400000):
    """柔和竖向渐变填充：从 top(RGB tuple) 到 bottom(RGB tuple)。"""
    spPr = shape._element.spPr
    for tag in ('a:solidFill', 'a:gradFill', 'a:noFill', 'a:blipFill', 'a:pattFill', 'a:grpFill'):
        for e in spPr.findall(qn(tag)):
            spPr.remove(e)
    gf = OxmlElement('a:gradFill')
    gsLst = OxmlElement('a:gsLst')
    gs1 = OxmlElement('a:gs'); gs1.set('pos', '0')
    c1 = OxmlElement('a:srgbClr'); c1.set('val', '%02X%02X%02X' % top); gs1.append(c1)
    gs2 = OxmlElement('a:gs'); gs2.set('pos', '100000')
    c2 = OxmlElement('a:srgbClr'); c2.set('val', '%02X%02X%02X' % bottom); gs2.append(c2)
    gsLst.append(gs1); gsLst.append(gs2)
    lin = OxmlElement('a:lin'); lin.set('ang', str(angle)); lin.set('scaled', '1')
    gf.append(gsLst); gf.append(lin)
    spPr.append(gf)

def _soft_shadow(shape, layers):
    """多层轻柔外阴影，颜色统一为极深蓝 (15,30,60)，不发黑发重。
    layers: list of (blur_pt, dist_pt, alpha_pct, (r,g,b))"""
    spPr = shape._element.spPr
    effectLst = spPr.find(qn('a:effectLst'))
    if effectLst is None:
        effectLst = OxmlElement('a:effectLst'); spPr.append(effectLst)
    for (blur, dist, alpha, col) in layers:
        sh = OxmlElement('a:outerShdw')
        sh.set('blur', str(int(blur * 12700)))
        sh.set('dist', str(int(dist * 12700)))
        sh.set('dir', '5400000')
        sh.set('rotWithShape', '0')
        c = OxmlElement('a:srgbClr'); c.set('val', '%02X%02X%02X' % col)
        a = OxmlElement('a:alpha'); a.set('val', str(int(alpha * 1000))); c.append(a)
        sh.append(c); effectLst.append(sh)

def _round_corners(shape, adj=30000):
    """大圆角：adj 越大越圆（默认约 16-18px 级别）。"""
    prstGeom = shape._element.spPr.find(qn('a:prstGeom'))
    if prstGeom is None:
        return
    avLst = prstGeom.find(qn('a:avLst'))
    if avLst is None:
        avLst = OxmlElement('a:avLst'); prstGeom.append(avLst)
    gd = avLst.find(qn('a:gd'))
    if gd is None:
        gd = OxmlElement('a:gd'); avLst.append(gd)
    gd.set('name', 'adj'); gd.set('fmla', 'val %d' % adj)

def card(slide, left, top, width, height, title, desc, action=None, tsize=12.5, dsize=10.5,
         title_font=FONT, title_color=NAVY, body_font=FONT, accent_token=False,
         action_lead_color=None, action_color=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    _gradient_fill(sh, (0xFF, 0xFF, 0xFF), (0xF4, 0xF8, 0xFD))
    sh.line.color.rgb = LIGHT_BORDER; sh.line.width = Pt(1)
    _round_corners(sh, 30000)
    _soft_shadow(sh, [(8, 3, 6, (15, 30, 60)), (14, 6, 8, (15, 30, 60)), (22, 9, 4, (15, 30, 60))])
    tf = sh.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.14); tf.margin_right = Inches(0.14); tf.margin_top = Inches(0.09); tf.margin_bottom = Inches(0.09)
    p = tf.paragraphs[0]; p.space_after = Pt(2)
    if accent_token:
        toks = title.split(' ')
        aidx = next((i for i, t in enumerate(toks) if t.startswith('卡点') and t[2:].isdigit()), -1)
        if aidx >= 0:
            for i, t in enumerate(toks):
                run = p.add_run(); run.text = ('' if i == 0 else ' ') + t
                run.font.size = Pt(tsize); run.font.bold = True
                run.font.color.rgb = TECH_BLUE if i == aidx else title_color; set_font(run, title_font)
        else:
            r = p.add_run(); r.text = title; r.font.size = Pt(tsize); r.font.bold = True; r.font.color.rgb = title_color; set_font(r, title_font)
    else:
        r = p.add_run(); r.text = title; r.font.size = Pt(tsize); r.font.bold = True; r.font.color.rgb = title_color; set_font(r, title_font)
    p2 = tf.add_paragraph(); p2.space_after = Pt(2 if action else 0); p2.line_spacing = 1.1
    r2 = p2.add_run(); r2.text = desc; r2.font.size = Pt(dsize); r2.font.color.rgb = GRAY_BODY; set_font(r2, body_font)
    if action:
        p3 = tf.add_paragraph(); p3.line_spacing = 1.1; p3.space_before = Pt(3)
        lead = action_lead_color if action_lead_color is not None else GREEN
        body = action_color if action_color is not None else GRAY_BODY
        r3 = p3.add_run(); r3.text = "对策："; r3.font.size = Pt(dsize); r3.font.color.rgb = lead; r3.font.bold = True; set_font(r3, body_font)
        r3b = p3.add_run(); r3b.text = action; r3b.font.size = Pt(dsize); r3b.font.color.rgb = body; set_font(r3b, body_font)
    return sh

def stat_card(slide, left, top, width, height, big, label, num_font=NUM_FONT, num_color=NAVY, light_shadow=False):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    _gradient_fill(sh, (0xFF, 0xFF, 0xFF), (0xEE, 0xF3, 0xFB))
    sh.line.color.rgb = LIGHT_BORDER_B; sh.line.width = Pt(1)
    _round_corners(sh, 30000)
    if light_shadow:
        _soft_shadow(sh, [(12, 3, 5, (15, 30, 60))])
    else:
        _soft_shadow(sh, [(8, 3, 6, (15, 30, 60)), (14, 6, 8, (15, 30, 60)), (22, 9, 4, (15, 30, 60))])
    tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = big; r.font.size = Pt(33); r.font.bold = True; r.font.color.rgb = num_color; set_font(r, num_font)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = label; r2.font.size = Pt(13); r2.font.color.rgb = GRAY_BODY; set_font(r2)
    # 左侧淡蓝细高亮线，突出数据感（SaaS 数据看板风格）
    hl = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Pt(5), Inches(height))
    _round_corners(hl, 60000)
    hl.fill.solid(); hl.fill.fore_color.rgb = ACCENT; hl.line.fill.background()
    return sh

TOTAL = 7

# ============ 1. COVER ============
s = prs.slides.add_slide(BLANK); add_bg(s, COVER_IMG)
text_box(s, 1.0, 2.0, 11.3, 1.2, "联邦智研 AISci", size=54, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font=SERIF_FONT)
gold_line(s, 5.2, 3.25, 2.9, 3)
text_box(s, 1.0, 3.5, 11.3, 0.7, "面向《Science》125 前沿科学问题的多智能体科研自动化系统", size=20, color=LIGHT_GOLD, align=PP_ALIGN.CENTER, font=SERIF_FONT)
text_box(s, 1.0, 4.45, 11.3, 0.6, "赛题指标达成与阶段进展报告", size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font=SERIF_FONT)
text_box(s, 1.0, 6.2, 11.3, 0.5, "团队：联邦智研团队　|　报告日期：2026-08-02", size=15, color=RGBColor(0xC8, 0xD2, 0xE6), align=PP_ALIGN.CENTER)

# ============ 2. (一)组队情况 + (二)a 指标达成（合并） ============
s = content_slide("组队与指标", "团队构成 + 核心指标高达成", 1, TOTAL, title_kind="academic", sub_color=RGBColor(0x64,0x74,0x8B), sub_font=SANS_FONT)
text_box(s, 0.5, 1.12, 5.6, 0.36, "（一）组队情况", size=15, color=NAVY, bold=True, font=SANS_FONT)
card(s, 0.5, 1.52, 5.6, 1.5, "1.1 团队构成", "总计 8 人（本科生 2 + 研究生 6），学科交叉互补，覆盖计算机/AI/数据科学，支撑多智能体与科研自动化全栈开发。", tsize=11, dsize=8.5, title_font=SANS_FONT, title_color=NAVY, body_font=SANS_FONT)
card(s, 0.5, 3.25, 5.6, 1.62, "1.2 主答辩人员", "梁亮雨（负责人，大三），获数学竞赛省二、交科赛国三等；当前短板：暂无答辩经验（已列入卡点对症施策）。", tsize=11, dsize=8.5, title_font=SANS_FONT, title_color=NAVY, body_font=SANS_FONT)
card(s, 0.5, 5.05, 5.6, 1.5, "1.3 开发投入", "核心开发人员 7.13–8.1 留校集中开发，已完成大部分工作，为 8 月中旬交付奠定坚实基础。", tsize=11, dsize=8.5, title_font=SANS_FONT, title_color=NAVY, body_font=SANS_FONT)
text_box(s, 6.5, 1.12, 6.3, 0.36, "（二）指标达成（总体完成率）", size=15, color=NAVY, bold=True, font=SANS_FONT)
stat_card(s, 6.5, 1.5, 2.9, 1.15, "8/8", "评分指标已达成", num_font=NUM_FONT, num_color=INK_BLUE, light_shadow=True)
stat_card(s, 9.6, 1.5, 2.9, 1.15, "100%", "完全达成率", num_font=NUM_FONT, num_color=INK_BLUE, light_shadow=True)
stat_card(s, 6.5, 2.78, 2.9, 1.15, "0", "未达成/待补强项", num_font=NUM_FONT, num_color=INK_BLUE, light_shadow=True)
stat_card(s, 9.6, 2.78, 2.9, 1.15, "3/3", "赛道主线 A/B/B′ 全达成", num_font=NUM_FONT, num_color=INK_BLUE, light_shadow=True)
bullets(s, 6.5, 4.15, 6.0, 2.4, [
    "8 项评分指标：已达成 8 / 未达成 0（无部分达成项）。",
    "指标完全达成率 = 8/8 = 100%。",
    "三项赛道主线 A / B / B′ 全部“已达成”。",
    "结论：闭环科研能力完整，核心指标达成度高，全部 8 项均已达成。",
], size=11.5, gap=6, font=SANS_FONT)

# ============ 4. (二)b 指标总览大表 ============
s = content_slide("指标总览", "9 项评分指标 + 3 赛道主线达成状态", 2, TOTAL, title_kind="academic", sub_color=RGBColor(0x64,0x74,0x8B), sub_font=SANS_FONT)
ind = [
    ["科学价值", "科学事实表达准确性", "15", "核心", "已达成", "Fact 白名单 + 证据链双向溯源；ReportQualityCheck 强制查重，参考文献零虚构（critical_issues 为空）"],
    ["科学价值", "内容转化/解释/展示清晰度", "15", "一般", "已达成", "证据链抽屉、PipelineProgress、RunLogDetail、QualityCheckCard、影响力雷达图等多层可视化"],
    ["科学价值", "作品主题完整性与一致性", "10", "一般", "已达成", "七阶段闭环逻辑一致，A/B/B′ 三大赛道方向均有产出"],
    ["技术深度", "模型/智能体/技能设计完整性", "10", "核心", "已达成", "6 个智能体（独立类+Prompt+Schema）、约 70 个技能、11 种阶段布尔门禁"],
    ["技术深度", "结果校验/反馈迭代/稳定性设计", "10", "核心", "已达成", "布尔门禁 + 停滞停止 + 闭环决策 + 审计链 + HITL + 科学自迭代编排"],
    ["应用潜力", "面向真实场景的使用价值", "10", "一般", "已达成", "高校科研训练营/毕设试点稳定运行；Science125/评分表/联邦学习仿真三类场景验证；HCOE 教育赋能模型落地"],
    ["应用潜力", "作品演示/交互入口/交付完整度", "10", "核心", "已达成", "前端 8 页面 + 预测 Tab + 一键启动 + Cloudflare 内网穿透 + 答辩 PPT"],
    ["应用潜力", "代码/结果/流程可复现性", "10", "核心", "已达成", "USE_MOCK_LLM + 审计链 jsonl 导出 + pytest + check_e2e（14 项）+ 一键脚本"],
    ["赛道 A", "科学假设生成与研究计划设计", "—", "核心", "已达成", "七阶段闭环，假设绑定 supporting_fact_ids，证据弱自动补文献重跑"],
    ["赛道 B", "科学实验任务规划与反馈迭代", "—", "核心", "已达成", "迭代实验 Tab：绑数据→设计脚本→沙箱执行→结果分析→反馈重设计"],
    ["赛道 B′", "科研影响力分析与偏差解释", "—", "核心", "已达成", "pingfenbiao D1–D4，D4 识别 7 类偏差并给出校准与缓解措施"],
]
nr, nc = len(ind) + 1, 6
tbl = s.shapes.add_table(nr, nc, Inches(0.4), Inches(1.15), Inches(12.5), Inches(5.55)).table
cw = [1.4, 2.2, 0.7, 0.95, 1.15, 6.1]
for i, w in enumerate(cw): tbl.columns[i].width = Inches(w)
hdr = ["维度", "考核指标", "满分", "类别", "达成状态", "关键证据"]
for ci, h in enumerate(hdr):
    c = tbl.cell(0, ci); c.fill.solid(); c.fill.fore_color.rgb = NAVY
    c.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = h; r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = WHITE; set_font(r)
for ri, row in enumerate(ind, start=1):
    for ci, val in enumerate(row):
        c = tbl.cell(ri, ci); c.fill.solid(); c.fill.fore_color.rgb = WHITE if ri % 2 else RGBColor(0xEE, 0xF2, 0xF8)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = Inches(0.05); c.margin_right = Inches(0.05); c.margin_top = Inches(0.02); c.margin_bottom = Inches(0.02)
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if ci in (1, 5) else PP_ALIGN.CENTER
        r = p.add_run(); r.text = val; r.font.size = Pt(9.5); r.font.name = FONT
        if ci == 4:
            r.font.color.rgb = GREEN if "已达成" in val else GOLD; r.font.bold = True
        else:
            r.font.color.rgb = DARK_TEXT
        set_font(r)

# ============ 4. (二)c 差距与区分 + (三) 单位对接（合并） ============
s = content_slide("区分与对接", "关键/一般指标区分 · 单位对接", 3, TOTAL, title_kind="academic", sub_color=RGBColor(0x64,0x74,0x8B), sub_font=SANS_FONT)
sep = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.1), Pt(2), Inches(5.6))
sep.fill.solid(); sep.fill.fore_color.rgb = GOLD; sep.line.fill.background()
# 左：差距与区分
text_box(s, 0.5, 1.1, 5.9, 0.34, "（二）关键/一般指标区分", size=16, color=NAVY, bold=True, font=SANS_FONT)
text_box(s, 0.5, 1.5, 5.9, 1.0, "2.3 区分说明：全部 8 项评分指标均已达成；核心指标决定获奖、评委重点关注，一般指标为加成项。", size=12.5, color=GRAY_BODY, font=SANS_FONT)
text_box(s, 0.5, 3.15, 2.9, 0.32, "【核心指标·全达成】", size=13, color=NAVY, bold=True, font=SANS_FONT)
bullets(s, 0.5, 3.5, 2.85, 3.1, [
    "科学事实表达准确性",
    "结果校验/反馈迭代/稳定性",
    "代码/结果可复现性",
    "模型/智能体/技能完整性",
    "赛道 A/B/B′ 主线达成",
    "演示/交付完整度",
], size=12, gap=2, color=GRAY_BODY, font=SANS_FONT)
text_box(s, 3.4, 3.15, 2.9, 0.32, "【一般指标】", size=13, color=NAVY, bold=True, font=SANS_FONT)
bullets(s, 3.4, 3.5, 2.85, 3.1, [
    "内容展示清晰度(已达成)",
    "主题完整一致性(已达成)",
    "真实场景使用价值(已达成)",
], size=12, gap=4, color=GRAY_BODY, font=SANS_FONT)
# 右：单位对接
text_box(s, 6.7, 1.1, 6.1, 0.34, "（三）单位对接", size=16, color=NAVY, bold=True, font=SANS_FONT)
bullets(s, 6.7, 1.5, 6.2, 4.0, [
    ("3.1 渠道与频次", "钉钉常态化答疑 + 每周直播双通道，高频保障及时响应。"),
    ("3.2 对接层级", "直达发榜单位技术负责人，链路短、无信息衰减。"),
    ("3.3 技术路线对齐", "启动即对齐路线/边界/期望，贴合评审导向。"),
    ("3.4 难点澄清", "“鼓励/必须/重点难点”三张清单，关键事项已闭环验证。"),
    ("3.5 资源支持", "提供 App 与有限大模型免费额度，降低试错成本。"),
], size=13.5, gap=8, color=GRAY_BODY, head_color=NAVY, font=SANS_FONT)
text_box(s, 6.7, 5.6, 6.2, 1.1, "小结：呈现“高频答疑 + 直达技术负责人”特点，在路线对齐、难点澄清与资源供给三方面形成有效支撑。", size=13.5, color=NAVY, bold=True, font=SANS_FONT)

# ============ 7. (四) 堵点卡点（技术/资源/协调政策 合并） ============
s = content_slide("堵点卡点", "4.1 技术难题 · 4.2 资源缺口 · 4.3/4.4 协调政策（共 9 项）", 4, TOTAL, title_kind="academic", sub_color=RGBColor(0x64,0x74,0x8B), sub_font=SANS_FONT)
col_defs = [(0.4, "【4.1 技术难题】"), (4.6, "【4.2 资源缺口】"), (8.8, "【4.3/4.4 协调政策】")]
tech = [
    ("卡点1 多模态证据链深度不足", "仅图像/图表进证据链，化学式/质谱/音视频缺统一表示与对齐。", "对策：化学式解析(RDKit)+多模态技能层扩展音视频与谱图。"),
    ("卡点2 反事实预演仅 L0 定性层", "仅做可证伪性过滤，未对实验成效做数值预演。", "对策：引入轻量数值代理模型(surrogate)做 L1 前评估。"),
    ("卡点3 长链路偶发“假设退化”", "极端证据缺失时迭代收益递减。", "对策：迭代收益<阈值自动冻结的强化规则。"),
]
reso = [
    ("卡点1 算力与 API 额度紧张", "Science125 全量跑批/联邦仿真额度紧，本地缺 GPU。", "对策：申请阿里云算力资助或教育优惠额度。"),
    ("卡点2 高质量标注数据稀缺", "跨校真实科研过程数据难获取，影响 D3 校准。", "对策：OpenAlex 元数据+平台自产审计链冷启动。"),
    ("卡点3 缺少 PPT 样例或模板", "答辩缺可参考样例，PPT 规范度受限，答辩人缺经验。", "对策：提炼统一模板+生成 PPT v1 骨架+加演练。"),
]
coord = [
    ("4.3 卡点1 论文/数据集访问受限", "OpenAlex/Zenodo 等公开数据源有限。", "对策：企业渠道建库+数据快速处理提示词。"),
    ("4.4 卡点1 AI 生成内容边界", "需明确 AI 生成内容的署名与责任边界。", "对策：强制标注 actual/simulated/expected+伦理说明。"),
    ("4.4 卡点2 数据合规与跨境", "跨校数据涉及个人信息保护。", "对策：仅用公开/脱敏数据+签署数据使用协议。"),
]
cols = [tech, reso, coord]
top0, ch = 1.62, 1.58
for ci, (cl, chead) in enumerate(col_defs):
    text_box(s, cl, 1.12, 4.0, 0.36, chead, size=13.5, color=NAVY, bold=True, font=SANS_FONT)
    for ri in range(3):
        t, d, a = cols[ci][ri]
        card(s, cl, top0 + ri * (ch + 0.16), 3.95, ch, t, d, action=a, tsize=11, dsize=9,
             title_font=SANS_FONT, title_color=NAVY, body_font=SANS_FONT,
             action_lead_color=GREEN, action_color=GRAY_BODY)

# ============ 10. (五) 工作计划表 ============
s = content_slide("工作计划", "四项里程碑，8/15 前交付闭环", 5, TOTAL, title_kind="academic", sub_color=RGBColor(0x64,0x74,0x8B), sub_font=SANS_FONT)
plan = [
    ["里程碑任务", "时间节点", "负责人"],
    ["反事实预演 L1 轻量数值代理模型（联邦学习场景先行）", "8/08 前", "张豪杰、赵豪杰"],
    ["可复现性加固：一键脚本跨平台验证 + 审计链导出样例", "8/12 前", "廖嘉仪、韩克润、黄飞润"],
    ["闭环稳定性强化：迭代收益冻结规则 + 停滞停止阈值调优", "8/13 前", "陶实际、曾恩赐"],
    ["竞赛交付物整合：答辩 PPT v13、阶段报告、Demo 视频、原创性声明", "8/15 前", "全体 + 指导教师"],
]
nr, nc = len(plan), 3
tbl = s.shapes.add_table(nr, nc, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.2)).table
tbl.columns[0].width = Inches(7.6); tbl.columns[1].width = Inches(1.9); tbl.columns[2].width = Inches(2.8)
for ri, row in enumerate(plan):
    for ci, val in enumerate(row):
        c = tbl.cell(ri, ci); c.fill.solid()
        c.fill.fore_color.rgb = NAVY if ri == 0 else (WHITE if ri % 2 else RGBColor(0xEE, 0xF2, 0xF8))
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = Inches(0.08); c.margin_right = Inches(0.08); c.margin_top = Inches(0.04); c.margin_bottom = Inches(0.04)
        p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
        r = p.add_run(); r.text = val; r.font.size = Pt(13 if ri else 14)
        r.font.color.rgb = WHITE if ri == 0 else DARK_TEXT; r.font.bold = (ri == 0); set_font(r)

# ============ 8. (六) 综合自评（潜力/优势/风险 合并） ============
s = content_slide("综合自评", "获奖潜力高 · 主要优势 · 最大风险与突围", 6, TOTAL, title_kind="academic", sub_color=RGBColor(0x64,0x74,0x8B), sub_font=SANS_FONT)
text_box(s, 0.5, 1.12, 12.3, 0.6, "6.1 冲刺获奖潜力：☑ 高　理由：核心指标全部达成、闭环链条完整、可复现性强、差异化创新明确，8 月中旬前可进一步收敛提升。", size=13.5, color=NAVY, bold=True)
text_box(s, 0.5, 1.92, 6.0, 0.4, "6.2 主要优势", size=15, color=NAVY, bold=True)
bullets(s, 0.5, 2.38, 6.1, 4.3, [
    "方法创新清晰：Fact 白名单杜绝幻觉、布尔门禁替代连续评分、反馈中心+HITL、反事实 L0，组合差异化明显。",
    "闭环完整可审计：七阶段 Pipeline+科学自迭代编排+全链路审计链 jsonl 导出。",
    "有效提交率高：以“有效提交率”评估生成报告竞争力，Science125 一次成功率 96%（120/125），理想状态可达 100%。",
    "可复现性工程扎实：一键启动、pytest/check_e2e、审计链导出，第三方可复现。",
    "国产大模型底座(Qwen)+教育赋能叙事(HCOE)贴合政策与赛事导向。",
    "三大赛道方向(A/B/B′)均有可演示成果，覆盖面广。",
], size=11.5, gap=5)
sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.92), Inches(5.95), Inches(4.8))
sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0xFB, 0xF3, 0xE6); sh.line.color.rgb = GOLD; sh.line.width = Pt(1.5)
tf = sh.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.22); tf.margin_right = Inches(0.22); tf.margin_top = Inches(0.18)
p = tf.paragraphs[0]; r = p.add_run(); r.text = "6.3 最大风险点"
r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = NAVY; set_font(r)
p2 = tf.add_paragraph(); p2.space_before = Pt(6); p2.line_spacing = 1.2
r2 = p2.add_run(); r2.text = ("赛题设多方向（发榜单位众多、赛道交叉），评审维度交叉、细则开放，存在“覆盖面广但不够聚焦”被稀释风险；"
    "评分标准偏原则性，若材料未精准命中评审关注点，易被海量作品淹没；参赛人数多、强队密集，头部竞争激烈。")
r2.font.size = Pt(13); r2.font.color.rgb = DARK_TEXT; set_font(r2)
p3 = tf.add_paragraph(); p3.space_before = Pt(10)
r3 = p3.add_run(); r3.text = "突围策略：以“可审计闭环+零幻觉引用+可复现”三张差异化名片精准命中评审关注点，并在答辩现场强化记忆点呈现，8 月中旬前进一步收敛提升。"
r3.font.size = Pt(12.5); r3.font.color.rgb = GREEN; r3.font.bold = True; set_font(r3)

prs.save(OUT_PPTX)
print("SAVED ->", OUT_PPTX, os.path.getsize(OUT_PPTX), "bytes; slides:", len(prs.slides._sldIdLst))
