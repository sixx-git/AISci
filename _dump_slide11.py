from pptx import Presentation
from pptx.shapes.group import GroupShape

def dump_shape(shp, indent=0):
    sp = '  '*indent
    extra = ''
    if shp.shape_type == 13:
        extra = ' [PICTURE]'
    elif shp.shape_type == 3:
        extra = ' [CHART]'
    print(f"{sp}{shp.name} ({shp.shape_type}) pos=({shp.left.inches:.2f},{shp.top.inches:.2f}) size=({shp.width.inches:.2f},{shp.height.inches:.2f}){extra}")
    if shp.has_text_frame:
        for j, para in enumerate(shp.text_frame.paragraphs):
            txt = ''.join([r.text for r in para.runs]).strip()
            if txt:
                print(f"{sp}  txt: {txt[:120]}")
    if isinstance(shp, GroupShape):
        for child in shp.shapes:
            dump_shape(child, indent+2)

p = r"D:/Workplace/AISci/output/蓝白色医疗器械--淘宝店：烦烦优创馆.pptx"
prs = Presentation(p)
print(f"=== Slide 11 (index 10) ===")
for i, shp in enumerate(prs.slides[10].shapes):
    print(f"\n--- top {i} ---")
    dump_shape(shp)
